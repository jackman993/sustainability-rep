"""
TCFD 報告合併工具
將 TCFD 報告插入到 Environment 報告中
"""
from pptx import Presentation
from pathlib import Path
from typing import List
import logging
from .path_manager import get_tcfd_report_path, update_session_activity

logger = logging.getLogger(__name__)

def insert_tcfd_slides(env_prs: Presentation, insert_position: int = 5) -> Presentation:
    """
    將 TCFD 報告的所有幻燈片插入到 Environment 報告中
    
    Args:
        env_prs: Environment 報告的 Presentation 對象
        insert_position: 插入位置（頁碼，從 1 開始，第 5 頁）
    
    Returns:
        合併後的 Presentation 對象
    """
    # 獲取 TCFD 報告路徑
    tcfd_path = get_tcfd_report_path()
    if not tcfd_path or not tcfd_path.exists():
        raise FileNotFoundError(f"找不到 TCFD 報告: {tcfd_path}")
    
    logger.info(f"讀取 TCFD 報告: {tcfd_path}")
    
    # 讀取 TCFD 報告
    tcfd_prs = Presentation(str(tcfd_path))
    
    # 更新活動時間
    update_session_activity()
    
    # 獲取 TCFD 的所有幻燈片
    tcfd_slides = list(tcfd_prs.slides)
    
    if not tcfd_slides:
        logger.warning("TCFD 報告為空，跳過插入")
        return env_prs
    
    logger.info(f"準備插入 {len(tcfd_slides)} 頁 TCFD 報告到第 {insert_position} 頁位置（應為 7 頁）")
    
    # 驗證 TCFD 報告頁數
    if len(tcfd_slides) != 7:
        logger.warning(f"TCFD 報告頁數為 {len(tcfd_slides)} 頁，預期為 7 頁")
    
    return _merge_presentations(env_prs, tcfd_prs, insert_position)

def _merge_presentations(env_prs: Presentation, tcfd_prs: Presentation, insert_pos: int) -> Presentation:
    """
    合併兩個 Presentation（實際實現）
    
    由於 python-pptx 的限制，我們需要創建新的 Presentation 並複製內容
    使用 XML 層面的複製以確保完整保留格式和內容
    """
    # 創建新的 Presentation（使用 Environment 報告的模板）
    merged_prs = Presentation()
    
    # 1. 複製 Environment 報告中插入位置之前的幻燈片（第 1-4 頁）
    for i in range(min(insert_pos - 1, len(env_prs.slides))):
        env_slide = env_prs.slides[i]
        _copy_slide_xml(env_slide, merged_prs)
    
    # 2. 添加 TCFD 的所有幻燈片（第 5-11 頁，7 頁）
    for tcfd_slide in tcfd_prs.slides:
        _copy_slide_xml(tcfd_slide, merged_prs)
    
    # 3. 複製 Environment 報告中插入位置之後的幻燈片（如果有，第 12 頁開始）
    for i in range(insert_pos - 1, len(env_prs.slides)):
        env_slide = env_prs.slides[i]
        _copy_slide_xml(env_slide, merged_prs)
    
    return merged_prs

def _copy_slide_xml(source_slide, target_prs: Presentation):
    """
    通過複製 XML 來複製整個幻燈片（保留所有格式和內容）
    
    使用更安全的方法：創建新幻燈片後，只複製形狀內容，避免佈局文件重複
    """
    from copy import deepcopy
    
    # 創建新幻燈片（使用源幻燈片的佈局）
    new_slide = target_prs.slides.add_slide(source_slide.slide_layout)
    
    # 清除新幻燈片的默認佔位符內容
    shapes_to_remove = []
    for shape in new_slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # 複製源幻燈片的所有形狀
    for shape in source_slide.shapes:
        # 跳過佔位符（已經通過佈局處理）
        if shape.is_placeholder:
            continue
            
        try:
            # 深拷貝形狀的 XML 元素
            shape_xml_copy = deepcopy(shape.element)
            # 添加到新幻燈片的形狀樹
            new_slide.shapes._spTree.append(shape_xml_copy)
        except Exception as e:
            # 如果 XML 複製失敗，嘗試手動複製
            logger.warning(f"XML 複製失敗，使用手動複製: {e}")
            _copy_slide_content_manual(source_slide, new_slide, shape)

def _copy_slide_content_manual(source_slide, target_slide, shape):
    """
    手動複製單個形狀（當 XML 複製失敗時的備用方案）
    
    注意：這是簡化實現，可能無法完全保留所有格式
    """
    try:
        # 獲取形狀的位置和大小
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # 根據形狀類型複製
        if hasattr(shape, 'text') and shape.text:
            # 文本框
            try:
                textbox = target_slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.clear()
                
                # 複製文本和格式
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        p = text_frame.add_paragraph()
                        p.text = paragraph.text
                        if paragraph.font:
                            p.font.name = paragraph.font.name
                            p.font.size = paragraph.font.size
                            p.font.bold = paragraph.font.bold
            except Exception as e:
                logger.warning(f"複製文本框失敗: {e}")
        
        elif hasattr(shape, 'image'):
            # 圖片
            try:
                image = shape.image
                image_bytes = image.blob
                target_slide.shapes.add_picture(image_bytes, left, top, width, height)
            except Exception as e:
                logger.warning(f"複製圖片失敗: {e}")
        
        elif shape.shape_type == 19:  # 表格
            # 表格需要特殊處理
            logger.warning("表格複製需要特殊處理，當前跳過")
        
    except Exception as e:
        logger.warning(f"手動複製形狀時出錯: {e}")

