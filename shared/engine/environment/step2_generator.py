"""
Environment Step 2 報告生成引擎
生成 12 頁 Environment 報告：
- 第 1 頁：Cover 頁
- 第 2 頁：環境政策
- 第 3 頁：環境管理架構
- 第 4 頁：環境目標與指標
- 第 5-11 頁：TCFD 報告（7 頁，插入）
- 第 12 頁：總結/附錄
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import streamlit as st
import logging
import os

from ..path_manager import get_environment_output_path, update_session_activity
from ..tcfd_merger import insert_tcfd_slides
from ..session_cleanup import smart_cleanup

logger = logging.getLogger(__name__)

def generate_environment_report() -> Path:
    """
    生成 Environment 報告（12 頁）
    
    Returns:
        生成的報告文件路徑
    """
    # 清理過期會話
    smart_cleanup()
    
    logger.info("開始生成 Environment 報告")
    
    # 創建新的 Presentation
    prs = Presentation()
    
    # 第 1 頁：Cover 頁
    logger.info("生成第 1 頁：Cover 頁")
    add_cover_page(prs)
    
    # 第 2 頁：環境政策
    logger.info("生成第 2 頁：環境政策")
    add_environment_policy_page(prs)
    
    # 第 3 頁：環境管理架構
    logger.info("生成第 3 頁：環境管理架構")
    add_environment_management_page(prs)
    
    # 第 4 頁：環境目標與指標
    logger.info("生成第 4 頁：環境目標與指標")
    add_environment_targets_page(prs)
    
    # 第 5-11 頁：插入 TCFD 報告（7 頁）
    logger.info("插入 TCFD 報告（第 5-11 頁，7 頁）")
    prs = insert_tcfd_slides(prs, insert_position=5)
    
    # 第 12 頁：總結/附錄（在 TCFD 插入後添加）
    logger.info("生成第 12 頁：總結/附錄")
    add_summary_page(prs)
    
    # 更新活動時間
    update_session_activity()
    
    # 保存報告 - 使用強制寫入機制（與 TCFD 一致）
    print("[DEBUG] ========== Environment 文件保存開始 ==========")
    
    # Step 1: 獲取輸出路徑
    try:
        print("[DEBUG] Step 1: 獲取輸出路徑...")
        output_path = get_environment_output_path()
        print(f"[DEBUG] Step 1 成功: {output_path}")
    except Exception as path_error:
        error_msg = f"[ERROR] Failed to get output path: {str(path_error)}"
        print(error_msg)
        import traceback
        full_traceback = traceback.format_exc()
        print(full_traceback)
        try:
            st.error(f"❌ 路徑獲取失敗: {str(path_error)}")
            with st.expander("詳細錯誤信息", expanded=True):
                st.code(full_traceback)
        except:
            pass
        raise Exception(error_msg) from path_error
    
    # Step 2: 確保目錄存在（強制創建）
    try:
        print(f"[DEBUG] Step 2: 強制創建目錄 {output_path.parent}...")
        
        # 方法 1: 標準 mkdir
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
        except Exception as e1:
            print(f"[WARNING] mkdir() 失敗: {e1}")
            # 方法 2: 使用 os.makedirs 強制創建
            try:
                os.makedirs(str(output_path.parent), exist_ok=True, mode=0o777)
                print(f"[DEBUG] os.makedirs() 成功")
            except Exception as e2:
                print(f"[ERROR] os.makedirs() 也失敗: {e2}")
                raise
        
        # 驗證目錄是否真的存在
        if not output_path.parent.exists():
            raise Exception(f"目錄創建失敗，目錄不存在: {output_path.parent}")
        
        # 嘗試修改權限確保可寫
        try:
            os.chmod(str(output_path.parent), 0o777)
        except:
            pass  # 權限修改失敗不影響
        
        print(f"[DEBUG] Step 2 成功: 目錄已創建/存在")
        print(f"[DEBUG] 目錄是否存在: {output_path.parent.exists()}")
        print(f"[DEBUG] 目錄是否可寫: {os.access(output_path.parent, os.W_OK)}")
    except Exception as dir_error:
        error_msg = f"[ERROR] Failed to create output directory: {str(dir_error)}"
        print(error_msg)
        import traceback
        full_traceback = traceback.format_exc()
        print(full_traceback)
        try:
            st.error(f"❌ 目錄創建失敗: {str(dir_error)}")
            with st.expander("詳細錯誤信息", expanded=True):
                st.code(full_traceback)
        except:
            pass
        raise Exception(error_msg) from dir_error
    
    # Step 3: 保存文件（強制寫入，多次嘗試）
    try:
        print(f"[DEBUG] Step 3: 保存文件到 {output_path}...")
        print(f"[DEBUG] Presentation 對象: {type(prs)}")
        print(f"[DEBUG] Slides 數量: {len(prs.slides)}")
        
        # 確保路徑是字符串
        save_path = str(output_path)
        print(f"[DEBUG] 保存路徑 (字符串): {save_path}")
        
        # 強制保存文件（多次嘗試）
        print(f"[DEBUG] 準備保存文件到: {save_path}")
        
        # 確保父目錄存在且可寫（多次嘗試）
        parent_dir = Path(save_path).parent
        if not parent_dir.exists():
            print(f"[WARNING] 父目錄不存在，強制創建: {parent_dir}")
            # 方法 1: Path.mkdir
            try:
                parent_dir.mkdir(parents=True, exist_ok=True)
            except Exception as e1:
                print(f"[WARNING] Path.mkdir 失敗: {e1}")
                # 方法 2: os.makedirs
                try:
                    os.makedirs(str(parent_dir), exist_ok=True, mode=0o777)
                except Exception as e2:
                    print(f"[ERROR] os.makedirs 也失敗: {e2}")
                    raise Exception(f"無法創建父目錄: {e1}, {e2}")
        
        # 驗證父目錄存在
        if not parent_dir.exists():
            raise Exception(f"父目錄創建失敗，目錄不存在: {parent_dir}")
        
        # 檢查並修復權限
        if not os.access(parent_dir, os.W_OK):
            print(f"[WARNING] 父目錄不可寫，嘗試修改權限...")
            try:
                os.chmod(str(parent_dir), 0o777)
                print(f"[DEBUG] 權限修改完成")
            except Exception as perm_ex:
                print(f"[WARNING] 無法修改權限: {perm_ex}")
                # 即使權限修改失敗，也繼續嘗試保存
        
        print(f"[DEBUG] 父目錄狀態: 存在={parent_dir.exists()}, 可寫={os.access(parent_dir, os.W_OK)}")
        
        # 強制保存文件（多次嘗試不同方法）
        save_success = False
        last_error = None
        
        # 嘗試 1: 直接保存
        try:
            print(f"[DEBUG] 嘗試 1: 直接保存到 {save_path}")
            prs.save(save_path)
            print(f"[DEBUG] prs.save() 調用完成")
            save_success = True
        except Exception as save_ex1:
            print(f"[WARNING] 嘗試 1 失敗: {save_ex1}")
            last_error = save_ex1
            
            # 嘗試 2: 使用絕對路徑
            try:
                abs_path = Path(save_path).resolve()
                print(f"[DEBUG] 嘗試 2: 使用絕對路徑 {abs_path}")
                prs.save(str(abs_path))
                save_path = str(abs_path)  # 更新路徑
                print(f"[DEBUG] 使用絕對路徑保存成功")
                save_success = True
            except Exception as save_ex2:
                print(f"[WARNING] 嘗試 2 也失敗: {save_ex2}")
                last_error = save_ex2
                
                # 嘗試 3: 先刪除舊文件（如果存在）再保存
                try:
                    if Path(save_path).exists():
                        print(f"[DEBUG] 嘗試 3: 刪除舊文件後保存")
                        Path(save_path).unlink()
                    prs.save(save_path)
                    print(f"[DEBUG] 刪除舊文件後保存成功")
                    save_success = True
                except Exception as save_ex3:
                    print(f"[ERROR] 嘗試 3 也失敗: {save_ex3}")
                    last_error = save_ex3
        
        if not save_success:
            raise Exception(f"所有保存方法都失敗。最後錯誤: {last_error}")
        
        # 立即驗證文件是否真的存在
        import time
        time.sleep(0.2)  # 等待文件系統更新
        
        # 檢查多個可能的路徑
        check_paths = [
            Path(save_path),
            Path(save_path).resolve(),
        ]
        
        file_exists = False
        actual_path = None
        for check_path in check_paths:
            if check_path.exists():
                file_exists = True
                actual_path = check_path
                break
        
        if not file_exists:
            error_msg = f"[ERROR] 文件保存後不存在！保存路徑: {save_path}"
            print(error_msg)
            print(f"[DEBUG] 父目錄是否存在: {Path(save_path).parent.exists()}")
            print(f"[DEBUG] 父目錄: {Path(save_path).parent}")
            print(f"[DEBUG] 父目錄是否可寫: {os.access(Path(save_path).parent, os.W_OK)}")
            # 列出父目錄中的所有文件
            try:
                files_in_parent = list(Path(save_path).parent.iterdir())
                print(f"[DEBUG] 父目錄中的文件: {files_in_parent}")
            except:
                pass
            raise Exception(error_msg)
        
        # 使用實際存在的路徑
        save_path = str(actual_path)
        print(f"[DEBUG] 文件確認存在於: {save_path}")
        
        file_size = Path(save_path).stat().st_size
        print(f"[DEBUG] Step 3 成功: 文件已保存")
        print(f"[DEBUG] 文件大小: {file_size} bytes")
        
        # 驗證文件大小是否合理（至少應該有幾 KB）
        if file_size < 1000:  # 小於 1KB 可能不正常
            print(f"[WARNING] 文件大小異常小: {file_size} bytes，可能保存不完整")
            try:
                st.warning(f"⚠️ 文件大小異常小 ({file_size} bytes)，可能保存不完整")
            except:
                pass
        
        # 更新最終路徑
        output_path = Path(save_path)
        logger.info(f"Environment 報告已保存: {output_path}")
        
        # 更新 session_state
        st.session_state["environment_report_file"] = str(output_path)
        
        print("[DEBUG] ========== Environment 文件保存完成 ==========")
        return output_path
        
    except Exception as save_error:
        error_msg = f"[ERROR] Failed to save PPTX file: {str(save_error)}"
        print(error_msg)
        import traceback
        full_traceback = traceback.format_exc()
        print(full_traceback)
        try:
            st.error(f"❌ Environment 報告保存失敗: {str(save_error)}")
            with st.expander("詳細錯誤信息", expanded=True):
                st.code(full_traceback)
        except:
            pass
        raise Exception(error_msg) from save_error

def add_cover_page(prs: Presentation):
    """添加 Cover 頁（第 1 頁）"""
    # 使用標題幻燈片佈局
    slide_layout = prs.slide_layouts[0]  # 標題幻燈片
    slide = prs.slides.add_slide(slide_layout)
    
    # 設置標題
    title = slide.shapes.title
    title.text = "環境報告"
    
    # 設置副標題（如果存在）
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Environment Report"
    
    logger.debug("Cover 頁已添加")

def add_environment_policy_page(prs: Presentation):
    """添加環境政策頁（第 2 頁）"""
    slide_layout = prs.slide_layouts[1]  # 標題和內容
    slide = prs.slides.add_slide(slide_layout)
    
    # 設置標題
    title = slide.shapes.title
    title.text = "環境政策"
    
    # 添加內容
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = "環境政策內容..."
        # TODO: 根據實際需求填充內容
    
    logger.debug("環境政策頁已添加")

def add_environment_management_page(prs: Presentation):
    """添加環境管理架構頁（第 3 頁）"""
    slide_layout = prs.slide_layouts[1]  # 標題和內容
    slide = prs.slides.add_slide(slide_layout)
    
    # 設置標題
    title = slide.shapes.title
    title.text = "環境管理架構"
    
    # 添加內容
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = "環境管理架構內容..."
        # TODO: 根據實際需求填充內容
    
    logger.debug("環境管理架構頁已添加")

def add_environment_targets_page(prs: Presentation):
    """添加環境目標與指標頁（第 4 頁）"""
    slide_layout = prs.slide_layouts[1]  # 標題和內容
    slide = prs.slides.add_slide(slide_layout)
    
    # 設置標題
    title = slide.shapes.title
    title.text = "環境目標與指標"
    
    # 添加內容
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = "環境目標與指標內容..."
        # TODO: 根據實際需求填充內容
    
    logger.debug("環境目標與指標頁已添加")

def add_summary_page(prs: Presentation):
    """添加總結/附錄頁（第 12 頁）"""
    slide_layout = prs.slide_layouts[1]  # 標題和內容
    slide = prs.slides.add_slide(slide_layout)
    
    # 設置標題
    title = slide.shapes.title
    title.text = "總結與附錄"
    
    # 添加內容
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = "總結與附錄內容..."
        # TODO: 根據實際需求填充內容
    
    logger.debug("總結/附錄頁已添加")

