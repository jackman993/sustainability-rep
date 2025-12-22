"""
分析原模板的母片結構和裝飾元素
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

template_path = r"C:\Users\User\Downloads\母片模板3.pptx"

# 如果檔案不存在，嘗試其他路徑
if not Path(template_path).exists():
    # 嘗試其他可能的路徑
    alt_paths = [
        r"C:\Users\User\Desktop\母片模板3.pptx",
        r"C:\Users\User\Downloads\templet_resaved.pptx",
    ]
    for alt_path in alt_paths:
        if Path(alt_path).exists():
            template_path = alt_path
            break

def analyze_template(template_path):
    """分析模板的母片結構"""
    try:
        prs = Presentation(template_path)
        
        print("="*60)
        print(f"分析模板：{template_path}")
        print("="*60)
        
        # 投影片尺寸
        print(f"\n📐 投影片尺寸：")
        print(f"  寬度：{prs.slide_width.inches:.2f}\" ({prs.slide_width.emu / 914400:.2f} cm)")
        print(f"  高度：{prs.slide_height.inches:.2f}\" ({prs.slide_height.emu / 914400:.2f} cm)")
        
        # 分析母片
        slide_master = prs.slide_masters[0]
        print(f"\n🎨 母片分析：")
        print(f"  版面配置數量：{len(slide_master.slide_layouts)}")
        
        # 分析每個版面配置
        for idx, layout in enumerate(slide_master.slide_layouts):
            print(f"\n  📄 版面配置 {idx}: {layout.name}")
            print(f"    形狀數量：{len(layout.shapes)}")
            
            for shape_idx, shape in enumerate(layout.shapes):
                shape_type = type(shape).__name__
                shape_name = getattr(shape, 'name', '無名稱')
                
                # 位置資訊
                left = shape.left.inches if hasattr(shape.left, 'inches') else shape.left / 914400
                top = shape.top.inches if hasattr(shape.top, 'inches') else shape.top / 914400
                width = shape.width.inches if hasattr(shape.width, 'inches') else shape.width / 914400
                height = shape.height.inches if hasattr(shape.height, 'inches') else shape.height / 914400
                
                print(f"\n    形狀 {shape_idx}: {shape_name}")
                print(f"      類型：{shape_type}")
                print(f"      位置：left={left:.2f}\", top={top:.2f}\"")
                print(f"      尺寸：width={width:.2f}\", height={height:.2f}\"")
                
                # 如果是文字框，顯示內容
                if hasattr(shape, 'text'):
                    text = shape.text[:50] if shape.text else "(空)"
                    print(f"      文字：{text}...")
                
                # 如果是形狀，顯示形狀類型
                if hasattr(shape, 'shape_type'):
                    print(f"      形狀類型：{shape.shape_type}")
                
                # 檢查是否為半圓形或弧形
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type in [MSO_SHAPE.ARC, MSO_SHAPE.OVAL, MSO_SHAPE.ROUND_RECTANGLE]:
                        print(f"      ⚠️ 可能是裝飾圖案（弧形/圓形）")
        
        # 分析母片本身的形狀
        print(f"\n🎨 母片本身形狀：")
        print(f"  形狀數量：{len(slide_master.shapes)}")
        
        for shape_idx, shape in enumerate(slide_master.shapes):
            shape_type = type(shape).__name__
            shape_name = getattr(shape, 'name', '無名稱')
            
            left = shape.left.inches if hasattr(shape.left, 'inches') else shape.left / 914400
            top = shape.top.inches if hasattr(shape.top, 'inches') else shape.top / 914400
            width = shape.width.inches if hasattr(shape.width, 'inches') else shape.width / 914400
            height = shape.height.inches if hasattr(shape.height, 'inches') else shape.height / 914400
            
            print(f"\n  形狀 {shape_idx}: {shape_name}")
            print(f"    類型：{shape_type}")
            print(f"    位置：left={left:.2f}\", top={top:.2f}\"")
            print(f"    尺寸：width={width:.2f}\", height={height:.2f}\"")
            
            # 檢查是否在右上角（right > 70% of width）
            slide_width_inches = prs.slide_width.inches
            if left > slide_width_inches * 0.7:
                print(f"    ⚠️ 位於右上角區域（可能是半圓形裝飾）")
            
            if hasattr(shape, 'shape_type'):
                print(f"    形狀類型：{shape.shape_type}")
        
        print("\n" + "="*60)
        print("分析完成")
        print("="*60)
        
    except Exception as e:
        print(f"❌ 分析失敗：{e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    analyze_template(template_path)

