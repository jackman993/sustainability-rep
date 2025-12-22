"""
ESG Report Generator - Environment Chapter PPTX Engine
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os
import glob
import copy
from lxml import etree

from config import ENVIRONMENT_CONFIG, ENVIRONMENT_IMAGE_MAPPING, TCFD_TABLES, ASSETS_PATH
from content_engine import ContentEngine

# 加入 assets 路徑
import sys
sys.path.insert(0, str(Path(__file__).parent / "assets"))
# Note: TCFD_main_pptx no longer needed - we insert TCFD PPTX file directly
try:
    from emission_pptx import create_emission_table_on_slide_right
except ImportError:
    print("  ⚠ Warning: emission_pptx module not found, emission table may not work")
    create_emission_table_on_slide_right = None
sys.path.append(ASSETS_PATH)

# ============ SASB 產業映射 ============
SASB_MAP = {
    # 采矿冶金
    "钢铁": ("EM-IS", "钢铁生产"), "水泥": ("EM-CM", "建材"), "煤矿": ("EM-CO", "煤炭"),
    "采矿": ("EM-MM", "金属矿业"), "石油": ("EM-OG", "油气开采"),
    
    # 能源公用
    "电力": ("IF-EU", "电力公用事业"), "天然气": ("IF-GU", "燃气"), "自来水": ("IF-WU", "水务"),
    
    # 运输
    "航空": ("TR-AF", "航空"), "物流": ("TR-AL", "空运物流"), "铁路": ("TR-RA", "铁路"),
    "货运": ("TR-RO", "公路运输"), "海运": ("TR-MT", "海运"), "租车": ("TR-CR", "租车"),
    
    # 房地产建筑
    "房地产": ("IF-RE", "不动产"), "建商": ("IF-HB", "建筑商"), "工程": ("IF-EN", "工程建设"),
    
    # 食品饮料
    "农业": ("FB-AB", "农产品"), "畜牧": ("FB-MP", "畜牧乳品"), "食品": ("FB-PF", "食品加工"),
    "餐饮": ("FB-RN", "餐饮"), "超市": ("FB-FR", "食品零售"), "饮料": ("FB-NB", "饮料"),
    
    # 医疗
    "制药": ("HC-BI", "生技制药"), "医材": ("HC-MS", "医疗器材"), "药局": ("HC-DY", "药品零售"),
    "医院": ("HC-DV", "医疗服务"),
    
    # 金融
    "银行": ("FN-CB", "商业银行"), "证券": ("FN-IB", "投资银行"), "保险": ("FN-IN", "保险"),
    "基金": ("FN-AC", "资产管理"), "消金": ("FN-CF", "消费金融"),
    
    # 科技
    "半导体": ("TC-SC", "半导体"), "电子": ("TC-HW", "硬件"), "软件": ("TC-SI", "软件IT"),
    "电信": ("TC-TL", "电信"), "电商": ("TC-ES", "电子商务"), "互联网": ("TC-IM", "网络媒体"),
    
    # 消费品
    "服装": ("CG-AA", "服饰"), "家电": ("CG-AM", "家电"), "家具": ("CG-BF", "家具"),
    "日用品": ("CG-HP", "家用品"), "零售": ("CG-MR", "零售"), "汽车": ("CG-AC", "汽车"),
    
    # 其他
    "化工": ("RT-CH", "化工"), "包装": ("RT-CP", "包装"), "机械": ("RT-IG", "工业机械"),
    "造纸": ("RR-PP", "纸浆造纸"), "太阳能": ("RR-ST", "太阳能"),
}

def get_sasb(industry_input):
    """根據產業名稱取得 SASB 代碼和名稱"""
    for keyword, (code, name) in SASB_MAP.items():
        if keyword in industry_input:
            return code, name
    return "通用", "一般产业"

# ============ 16:9 版面常數 ============
# 16:9 標準簡報尺寸: 10" x 5.625" (或 13.333" x 7.5")
# 使用標準簡報尺寸: 13.333" x 7.5"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 版面配置常數（16:9 比例）
MARGIN = Inches(0.4)                    # 頁面邊距
TITLE_TOP = Inches(0.3)                 # 標題頂部位置
TITLE_HEIGHT = Inches(0.7)              # 標題高度
CONTENT_TOP = Inches(1.1)               # 內容區頂部
CONTENT_HEIGHT = Inches(5.5)            # 內容區高度（16:9 調整）
CONTENT_WIDTH = Inches(6.0)             # 單側內容寬度（左或右，16:9 調整）
LEFT_CONTENT_LEFT = MARGIN              # 左側內容起點
RIGHT_CONTENT_LEFT = Inches(6.8)        # 右側內容起點（16:9 調整）
GAP = Inches(0.3)                       # 左右間距

# TCFD Generator 的輸出路徑
# Default path: C:\Users\User\Desktop\TCFD generator\output
# Alternative: C:\Users\User\Desktop\environment4.1-4.9\assets\TCFD
TCFD_OUTPUT_PATH = r"C:\Users\User\Desktop\TCFD generator\output"
TCFD_ASSETS_PATH = str(Path(__file__).parent / "assets" / "TCFD")  # Local assets folder

# Emission 輸出路徑
EMISSION_OUTPUT_PATH = Path(__file__).parent / "output"


class EnvironmentPPTXEngine:
    """Environment Chapter PPTX Report Generation Engine"""

    def __init__(self, template_path=None, test_mode=False, emission_data=None, industry="企業", tcfd_output_folder=None, emission_output_folder=None, company_profile=None, api_key=None):
        """
        Initialize engine
        template_path: Template file path (optional, defaults to handdrawppt.pptx in assets)
        test_mode: Test mode (skip LLM API)
        emission_data: Carbon emission data dict
        industry: Industry name
        tcfd_output_folder: TCFD output folder path (from Step 1)
        emission_output_folder: Emission output folder path (from Step 2)
        company_profile: Company size information dict (from Step 2)
        api_key: Claude API Key
        """
        self.emission_data = emission_data or {}
        self.industry = industry
        self.tcfd_output_folder = tcfd_output_folder  # Step 1 的 TCFD 輸出路徑
        self.emission_output_folder = emission_output_folder  # Step 2 的 Emission 輸出路徑
        self.company_profile = company_profile or {}
        self.api_key = api_key
        
        # Default template: handdrawppt.pptx in assets folder
        if template_path is None:
            default_template = Path(__file__).parent / "assets" / "handdrawppt.pptx"
            if default_template.exists():
                template_path = str(default_template)
                print(f"  ✓ Using default template: {template_path}")
            else:
                print(f"  ⚠ Default template not found: {default_template}")
        
        if template_path and os.path.exists(template_path):
            # Option C: Backup template to assets folder (using absolute path)
            import shutil
            env_report_dir = Path(__file__).parent
            template_backup_path = env_report_dir / "assets" / "template_backup.pptx"
            template_backup_path.parent.mkdir(exist_ok=True)  # Ensure folder exists
            if not template_backup_path.exists():
                shutil.copy(template_path, template_backup_path)
                print(f"✓ Template backed up to: {template_backup_path}")
            
            self.prs = Presentation(template_path)
            print(f"✓ Template loaded: {template_path}")
            
            # Delete template's default blank slides
            template_slide_count = len(self.prs.slides)
            if template_slide_count > 0:
                print(f"  ℹ Deleting template default slides ({template_slide_count} slides)")
                # Delete from back to front to avoid index issues
                for i in range(template_slide_count - 1, -1, -1):
                    rId = self.prs.slides._sldIdLst[i].rId
                    self.prs.part.drop_rel(rId)
                    del self.prs.slides._sldIdLst[i]
        else:
            self.prs = Presentation()
        
        # Set to 16:9 aspect ratio uniformly
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        print(f"✓ Slide size: 16:9 ({SLIDE_WIDTH.inches}\" x {SLIDE_HEIGHT.inches}\")")
        
        self.test_mode = test_mode
        
        # Debug: Check API Key
        if self.api_key:
            print(f"  ✓ EnvironmentPPTXEngine received API Key: {self.api_key[:10]}...")
        else:
            print(f"  ⚠ EnvironmentPPTXEngine did not receive API Key")
        
        self.content_engine = ContentEngine(
            test_mode=test_mode, 
            company_profile=self.company_profile,
            api_key=self.api_key
        )
        self.config = ENVIRONMENT_CONFIG
        
        # 版面設定
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        
        # 樣式設定
        self.title_font_size = Pt(28)
        self.body_font_size = Pt(12)  # Changed to 12pt
        self.font_name = 'Microsoft JhengHei'
        self.primary_color = RGBColor(26, 58, 46)  # 深綠色
        self.secondary_color = RGBColor(74, 124, 89)  # 淺綠色

    def _get_blank_layout(self):
        """Get blank layout"""
        # Prefer index 6 (standard blank layout)
        if len(self.prs.slide_layouts) > 6:
            return self.prs.slide_layouts[6]
        
        # Try to find layout with 'blank' in name
        for layout in self.prs.slide_layouts:
            if 'blank' in layout.name.lower() or layout.name == '空白':
                return layout
        
        # Finally use the first layout
        return self.prs.slide_layouts[0]

    def _add_slide(self):
        """Add blank slide (remove default placeholders and add watermark)"""
        layout = self._get_blank_layout()
        slide = self.prs.slides.add_slide(layout)
        
        # Remove layout's default placeholder shapes
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # Add watermark (bottom right corner)
        self._add_watermark(slide)
        
        return slide
    
    def _add_watermark(self, slide):
        """Add watermark to bottom right corner of slide"""
        try:
            # Calculate bottom right corner position
            watermark_width = Inches(2.5)
            watermark_height = Inches(0.3)
            watermark_left = SLIDE_WIDTH - watermark_width - MARGIN
            watermark_top = SLIDE_HEIGHT - watermark_height - Inches(0.2)
            
            # Add text box
            watermark_box = slide.shapes.add_textbox(watermark_left, watermark_top, watermark_width, watermark_height)
            watermark_frame = watermark_box.text_frame
            watermark_para = watermark_frame.paragraphs[0]
            watermark_para.text = "Sustainability Report"
            
            # Set format: italic, bold, 11pt, light gray
            watermark_para.font.size = Pt(11)
            watermark_para.font.bold = True
            watermark_para.font.italic = True
            watermark_para.font.name = 'Arial'
            watermark_para.font.color.rgb = RGBColor(200, 200, 200)  # Light gray
            watermark_para.alignment = PP_ALIGN.RIGHT
            
        except Exception as e:
            print(f"  ⚠ Watermark addition warning: {e}")

    def _add_title(self, slide, title_text, left=None, top=None, 
                   width=None, height=None):
        """Add title to slide (using 16:9 constants)"""
        left = left or MARGIN
        top = top or TITLE_TOP
        width = width or Inches(12.5)  # 16:9 width minus margins (13.333 - 0.4*2)
        height = height or TITLE_HEIGHT
        """Add title to slide"""
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = title_text
        p.font.size = self.title_font_size
        p.font.bold = True
        p.font.name = self.font_name
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.LEFT
        
        return title_box

    def _add_text_box(self, slide, text, left, top, width, height, 
                      font_size=None, bold=False, alignment=PP_ALIGN.LEFT):
        """Add text box to slide"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size or self.body_font_size
        p.font.bold = bold
        p.font.name = self.font_name
        p.alignment = alignment
        
        return text_box

    def _add_image(self, slide, image_name, left, top, width=None, height=None):
        """Add image to slide (from assets folder)"""
        image_path = os.path.join(ASSETS_PATH, image_name)
        return self._add_image_full_path(slide, image_path, left, top, width, height)

    def _add_image_full_path(self, slide, image_path, left, top, width=None, height=None):
        """Add image to slide (full path)"""
        if os.path.exists(image_path):
            try:
                if width and height:
                    pic = slide.shapes.add_picture(image_path, left, top, width, height)
                elif width:
                    pic = slide.shapes.add_picture(image_path, left, top, width=width)
                elif height:
                    pic = slide.shapes.add_picture(image_path, left, top, height=height)
                else:
                    pic = slide.shapes.add_picture(image_path, left, top)
                print(f"  ✓ Image inserted: {os.path.basename(image_path)}")
                return pic
            except Exception as e:
                print(f"  ✗ Image insertion failed: {image_path} - {e}")
        else:
            print(f"  ✗ Image not found: {image_path}")
        return None

    def _create_left_image_right_text_slide(self, title, image_name, text_content):
        """Create left image right text layout (16:9)"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 左側圖片
        self._add_image(slide, image_name, 
                       left=LEFT_CONTENT_LEFT, 
                       top=CONTENT_TOP, 
                       width=CONTENT_WIDTH,
                       height=CONTENT_HEIGHT)
        
        # 右側文字
        self._add_text_box(slide, text_content,
                          left=RIGHT_CONTENT_LEFT,
                          top=CONTENT_TOP,
                          width=CONTENT_WIDTH,
                          height=CONTENT_HEIGHT)
        
        return slide

    def _create_left_image_right_text_slide_full_path(self, title, image_path, text_content):
        """Create left image right text layout (16:9, image uses full path)"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 左側圖片
        self._add_image_full_path(slide, image_path, 
                                  left=LEFT_CONTENT_LEFT, 
                                  top=CONTENT_TOP, 
                                  width=CONTENT_WIDTH,
                                  height=CONTENT_HEIGHT)
        
        # 右側文字
        self._add_text_box(slide, text_content,
                          left=RIGHT_CONTENT_LEFT,
                          top=CONTENT_TOP,
                          width=CONTENT_WIDTH,
                          height=CONTENT_HEIGHT)
        
        return slide

    def _create_left_text_right_image_slide(self, title, text_content, image_name):
        """Create left text right image layout (16:9)"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 左側文字
        self._add_text_box(slide, text_content,
                          left=LEFT_CONTENT_LEFT,
                          top=CONTENT_TOP,
                          width=CONTENT_WIDTH,
                          height=CONTENT_HEIGHT)
        
        # 右側圖片
        self._add_image(slide, image_name,
                       left=RIGHT_CONTENT_LEFT,
                       top=CONTENT_TOP,
                       width=CONTENT_WIDTH,
                       height=CONTENT_HEIGHT)
        
        return slide

    def _create_full_text_slide(self, title, text_content):
        """Create text-only layout (16:9)"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 全版文字
        self._add_text_box(slide, text_content,
                          left=MARGIN,
                          top=CONTENT_TOP,
                          width=Inches(12.5),  # 16:9 寬度扣邊距 (13.333 - 0.4*2)
                          height=CONTENT_HEIGHT)
        
        return slide

    # ==================== 各章節生成方法 ====================

    def generate_cover_page(self):
        """Generate cover page"""
        print("\n[Generating Cover Page]")
        
        cover_text = self.content_engine.generate_environmental_cover(self.config)
        # Custom layout: left image 50% wider, right text 60% narrower (40% of original width)
        slide = self._add_slide()
        self._add_title(slide, "Chapter 4 Environmental Sustainability")
        
        # Left image: CONTENT_WIDTH * 1.5 = 9.0"
        left_image_width = CONTENT_WIDTH * 1.5
        self._add_image(slide, ENVIRONMENT_IMAGE_MAPPING['cover'],
                       left=LEFT_CONTENT_LEFT,
                       top=CONTENT_TOP,
                       width=left_image_width,
                       height=CONTENT_HEIGHT)
        
        # Right text: CONTENT_WIDTH * 0.4 = 2.4"
        right_text_width = CONTENT_WIDTH * 0.4
        right_text_left = LEFT_CONTENT_LEFT + left_image_width + GAP
        self._add_text_box(slide, cover_text,
                          left=right_text_left,
                          top=CONTENT_TOP,
                          width=right_text_width,
                          height=CONTENT_HEIGHT)
        
        print("✓ Cover page completed")

    def generate_policy_pages(self):
        """Generate 4.1-4.2 Environmental Policy Pages"""
        print("\n[Generating Environmental Policy Pages]")
        
        # Page 1: 4.1 Environmental Policy and Management Framework
        sustainability_text = self.content_engine.generate_sustainability_committee(self.config)
        self._create_left_image_right_text_slide(
            "4.1 Environmental Policy and Management Framework",
            ENVIRONMENT_IMAGE_MAPPING['sustainability_panel'],
            sustainability_text
        )
        
        # Page 2: 4.2 Four Dimensions of Environmental Policy
        policy_text = self.content_engine.generate_policy_description(self.config)
        self._create_left_text_right_image_slide(
            "4.2 Four Dimensions of Environmental Policy",
            policy_text,
            ENVIRONMENT_IMAGE_MAPPING['policy']
        )
        
        print("✓ Environmental policy pages completed")

    def _find_latest_tcfd_file(self):
        """Find the TCFD PPTX file (single file containing 7 pages)"""
        
        # TCFD file patterns (try multiple naming conventions)
        tcfd_patterns = [
            "TCFD*.pptx",            # General pattern (includes TCFD_table (26).pptx)
            "TCFD_*.pptx",           # Standard pattern
            "TCFD_7pages_*.pptx",    # Explicit 7 pages naming
            "TCFD_Complete_*.pptx",   # Complete TCFD file
        ]
        
        # Search paths in priority order
        search_paths = []
        
        # 1. Prefer using provided TCFD folder (from Step 1)
        if self.tcfd_output_folder and os.path.exists(self.tcfd_output_folder):
            search_paths.append(self.tcfd_output_folder)
            print(f"  ✓ Using Step 1 output folder: {self.tcfd_output_folder}")
        else:
            # 2. Check local assets folder first (for development/testing)
            if os.path.exists(TCFD_ASSETS_PATH):
                search_paths.append(TCFD_ASSETS_PATH)
                print(f"  ℹ Checking local assets folder: {TCFD_ASSETS_PATH}")
            
            # 3. Check default TCFD generator output path
            if os.path.exists(TCFD_OUTPUT_PATH):
                # Try to find latest subfolder
                tcfd_folders = glob.glob(os.path.join(TCFD_OUTPUT_PATH, "TCFD_*"))
                tcfd_folders = [f for f in tcfd_folders if os.path.isdir(f)]
                
                if tcfd_folders:
                    latest_folder = max(tcfd_folders, key=os.path.getmtime)
                    search_paths.append(latest_folder)
                    print(f"  ℹ Found latest TCFD folder: {os.path.basename(latest_folder)}")
                else:
                    search_paths.append(TCFD_OUTPUT_PATH)
                    print(f"  ℹ Checking default output folder: {TCFD_OUTPUT_PATH}")
        
        # Search in each path with each pattern
        for search_base in search_paths:
            for pattern in tcfd_patterns:
                search_path = os.path.join(search_base, pattern)
                files = glob.glob(search_path)
                if files:
                    # Get latest file
                    latest_file = max(files, key=os.path.getmtime)
                    print(f"  ✓ Found TCFD file: {os.path.basename(latest_file)}")
                    print(f"    Path: {latest_file}")
                    return latest_file
        
        print(f"  ✗ TCFD file not found in any search path")
        return None

    def _copy_slide_from_pptx(self, source_pptx_path):
        """Copy slide content from another PPTX"""
        try:
            source_prs = Presentation(source_pptx_path)
            
            for source_slide in source_prs.slides:
                # Add blank slide
                new_slide = self._add_slide()
                
                # Copy all shapes
                for shape in source_slide.shapes:
                    try:
                        # Copy shape XML
                        el = shape.element
                        new_el = copy.deepcopy(el)
                        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                    except Exception as e:
                        print(f"    ⚠ Warning when copying shape: {e}")
                
            return True
        except Exception as e:
            print(f"  ✗ Slide copy failed: {e}")
            return False

    def _insert_tcfd_pptx(self, pptx_path, title):
        """Insert all slides from TCFD PPTX file (should contain 7 pages)"""
        try:
            source_prs = Presentation(pptx_path)
            slide_count = len(source_prs.slides)
            
            print(f"  ℹ TCFD file contains {slide_count} slides")
            
            for idx, source_slide in enumerate(source_prs.slides):
                # Add blank slide (using template's white background)
                new_slide = self._add_slide()
                
                # Don't copy background, keep new slide's white background
                # Explicitly set background to white (ensure it's not black)
                try:
                    new_slide.background.fill.solid()
                    new_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                except:
                    pass
                
                # Copy all shapes (but not background)
                for shape in source_slide.shapes:
                    try:
                        el = shape.element
                        new_el = copy.deepcopy(el)
                        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                    except Exception as e:
                        # If copy fails, try other methods
                        pass
                
            print(f"  ✓ Inserted {slide_count} slides from {title}: {os.path.basename(pptx_path)}")
            return True
        except Exception as e:
            print(f"  ✗ Insertion failed {title}: {e}")
            return False

    def generate_tcfd_pages(self):
        """Generate TCFD pages (inserted from TCFD Generator - single PPTX file with 7 pages)"""
        print("\n[Generating TCFD Pages]")
        print("  Pages 5-11: Inserting TCFD PPTX file (7 pages)")
        
        # Find TCFD PPTX file (single file containing 7 pages)
        tcfd_file = self._find_latest_tcfd_file()
        
        if tcfd_file and os.path.exists(tcfd_file):
            # Insert all slides from the TCFD PPTX file (should be 7 pages)
            print(f"  ✓ Inserting TCFD slides from: {os.path.basename(tcfd_file)}")
            self._insert_tcfd_pptx(tcfd_file, "TCFD 7 Pages")
            
            # Verify slide count
            try:
                source_prs = Presentation(tcfd_file)
                slide_count = len(source_prs.slides)
                print(f"  ✓ TCFD file contains {slide_count} slides")
                if slide_count != 7:
                    print(f"  ⚠ Warning: Expected 7 slides, but found {slide_count} slides")
            except Exception as e:
                print(f"  ⚠ Could not verify slide count: {e}")
        else:
            # If file not found, create placeholder pages
            print(f"  ⚠ TCFD file not found, creating placeholder pages")
            placeholder_titles = [
                "TCFD Page 1",
                "TCFD Page 2",
                "TCFD Page 3",
                "TCFD Page 4",
                "TCFD Page 5",
                "TCFD Page 6",
                "TCFD Page 7"
            ]
            
            for title in placeholder_titles:
                slide = self._add_slide()
                self._add_title(slide, title)
                self._add_text_box(slide, 
                                 f"[{title}]\n\nPlease place TCFD PPTX file (7 pages) in:\n{TCFD_ASSETS_PATH}\n\nFile naming: TCFD*.pptx",
                                 left=Inches(0.5),
                                 top=Inches(1.5),
                                 width=Inches(12),
                                 height=Inches(5))
        
        print("✓ TCFD pages completed (Pages 5-11: 7 pages from single PPTX file)")

    def generate_sasb_page(self):
        """Generate SASB Industry Classification page (Page 12)"""
        print("\n[Generating SASB Page]")
        print("  Page 12: SASB Industry Classification")
        
        # SASB Industry Classification (inserted after TCFD, before GHG)
        sasb_slide = self._add_slide()
        self._add_title(sasb_slide, "SASB Industry Classification")
        
        # Get SASB code and name
        sasb_code, sasb_name = get_sasb(self.industry)
        
        # Left side: LLM-generated analysis text (150-170 words)
        sasb_analysis_text = self.content_engine.generate_sasb_analysis(
            self.config, 
            self.industry, 
            sasb_code, 
            sasb_name
        )
        self._add_text_box(sasb_slide, sasb_analysis_text,
                          left=LEFT_CONTENT_LEFT, 
                          top=CONTENT_TOP, 
                          width=CONTENT_WIDTH, 
                          height=CONTENT_HEIGHT,
                          font_size=Pt(12))
        
        # Right side: SASB table (16:9 half width, right-aligned)
        self._create_sasb_table(sasb_slide, sasb_code, sasb_name)
        
        print("✓ SASB page completed (Page 12)")

    def _create_sasb_table(self, slide, sasb_code, sasb_name):
        """Create SASB table on right half of slide (width reduced by 1.5cm, list all industry classifications)"""
        from pptx.oxml.xmlchemy import OxmlElement
        
        # 16:9 寬度
        slide_w = 13.333
        
        # 表格寬度：16:9 一半 - 1.5 公分 = 6.667" - 0.59" = 6.077"
        table_w = (slide_w / 2) - 0.59
        # 靠右對齊：slide_w - table_w - 右邊界(0.4)
        table_left = slide_w - table_w - 0.4
        
        # SASB 產業大類列表
        sasb_categories = [
            ("CG", "Consumer Goods", "消费品"),
            ("EM", "Extractives & Minerals Processing", "採掘與礦物加工"),
            ("FN", "Financials", "金融服務"),
            ("FB", "Food & Beverage", "食品飲料"),
            ("HC", "Health Care", "醫療保健"),
            ("IF", "Infrastructure", "基礎設施"),
            ("RR", "Renewable Resources & Alternative Energy", "可再生資源與替代能源"),
            ("RT", "Resource Transformation", "資源轉換（製造業）"),
            ("SV", "Services", "服務業"),
            ("TC", "Technology & Communications", "科技與通訊"),
            ("TR", "Transportation", "運輸"),
        ]
        
        # 取得當前產業的大類代碼（前綴）
        current_category_prefix = sasb_code.split("-")[0] if "-" in sasb_code else sasb_code
        
        # 表格位置：top 調整為 1.2" 以配合左側文字
        rows = 1 + len(sasb_categories)  # 表頭 + 11 個分類
        cols = 2  # 代碼、英文名稱（移除中文欄位）
        table_height = min(Inches(5.5), Inches(0.4 * rows))  # 動態高度，最多 5.5"
        table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(1.2), Inches(table_w), table_height)
        tbl = table_shape.table
        
        # 欄寬（調整為兩欄）
        tbl.columns[0].width = Inches(1.2)   # 代碼
        tbl.columns[1].width = Inches(4.877) # 英文名稱（佔用剩餘寬度）
        
        # 設定儲存格背景顏色
        def set_cell_fill(cell, hex_color):
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', hex_color)
            solidFill.append(srgbClr)
            tcPr.append(solidFill)
        
        # Header
        header_cell = tbl.cell(0, 0)
        header_cell.merge(tbl.cell(0, 1))
        header_cell.text = "SASB Industry Classification"
        set_cell_fill(header_cell, "2F5233")  # 深綠色
        p = header_cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
        header_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 填入所有產業分類（僅英文）
        for row_idx, (code, en_name, zh_name) in enumerate(sasb_categories, start=1):
            # 判斷是否為當前產業的分類
            is_current = (code == current_category_prefix)
            
            # 代碼欄
            code_cell = tbl.cell(row_idx, 0)
            code_cell.text = code
            if is_current:
                set_cell_fill(code_cell, "4A7C59")  # 深綠色（當前分類）
            else:
                set_cell_fill(code_cell, "F1F8F4")  # 極淺綠色
            p = code_cell.text_frame.paragraphs[0]
            if is_current:
                p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(10)
            p.font.bold = is_current
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER
            code_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 英文名稱欄（僅英文）
            en_cell = tbl.cell(row_idx, 1)
            en_cell.text = en_name
            if is_current:
                set_cell_fill(en_cell, "4A7C59")  # 深綠色（當前分類）
                p = en_cell.text_frame.paragraphs[0]
                p.font.color.rgb = RGBColor(255, 255, 255)
            else:
                set_cell_fill(en_cell, "F1F8F4")  # 極淺綠色
            p = en_cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.bold = is_current
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.LEFT
            en_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        print(f"  ✓ SASB table inserted (Code: {sasb_code}, Category: {sasb_name}, Current Classification: {current_category_prefix})")

    def _generate_emission_outputs(self):
        """Get emission outputs (prefer using Step 2 generated files)"""
        
        # Prefer using Step 2 provided output folder
        if self.emission_output_folder and os.path.exists(self.emission_output_folder):
            print(f"  ✓ Using Step 2 Emission output: {self.emission_output_folder}")
            
            # Find table and pie chart files
            table_files = glob.glob(os.path.join(self.emission_output_folder, "Emission_Table_*.pptx"))
            pie_files = glob.glob(os.path.join(self.emission_output_folder, "Emission_PieChart*.png"))
            
            results = {}
            if table_files:
                results["table_pptx"] = max(table_files, key=os.path.getmtime)
                print(f"  ✓ Found table: {os.path.basename(results['table_pptx'])}")
            if pie_files:
                results["pie_chart"] = max(pie_files, key=os.path.getmtime)
                print(f"  ✓ Found pie chart: {os.path.basename(results['pie_chart'])}")
            
            if results:
                return results
        
        # Fallback: call engine to regenerate
        try:
            print("  ℹ Step 2 output not found, regenerating...")
            from emission_pptx import generate_all
            results = generate_all()
            return results
        except Exception as e:
            print(f"  ⚠ Emission engine error: {e}")
            return None

    def generate_ghg_pages(self):
        """Generate Greenhouse Gas Emission Management Pages"""
        print("\n[Generating GHG Management Pages]")
        print("  Pages 13-15: GHG Management")
        
        # Set emission data from step1 emission engine BEFORE generating table
        # This ensures the table uses the calculated results from step1
        if self.emission_data:
            try:
                from assets.emission_pptx import set_emission_data
                set_emission_data(self.emission_data)
                print(f"  ✓ Emission data loaded from step1: Total {self.emission_data.get('total', 0):.2f} tCO₂e")
            except Exception as e:
                print(f"  ⚠ Unable to load emission data: {e}")
        
        # First generate emission engine outputs
        emission_results = self._generate_emission_outputs()
        
        # Page 13: 4.5 Carbon Inventory Table (left text right table)
        section_slide = self._add_slide()
        self._add_title(section_slide, "4.5 Carbon Inventory Table")
        
        # Left side: Text description
        ghg_text = self.content_engine.generate_ghg_calculation_method(self.config)
        self._add_text_box(section_slide, ghg_text,
                          left=LEFT_CONTENT_LEFT, 
                          top=CONTENT_TOP, 
                          width=CONTENT_WIDTH, 
                          height=CONTENT_HEIGHT,
                          font_size=Pt(12))
        
        # Right side: Emission table (reduced 50%, right-aligned)
        # This will use the emission_data set above from step1
        create_emission_table_on_slide_right(section_slide)
        
        # Page 14: Electricity Usage and Energy Conservation Policy (using emission pie chart)
        electricity_text = self.content_engine.generate_electricity_policy(self.config)
        if emission_results and "pie_chart" in emission_results:
            # Use emission-generated pie chart
            pie_chart_path = emission_results["pie_chart"]
            self._create_left_image_right_text_slide_full_path(
                "Electricity Usage and Energy Conservation Policy",
                str(pie_chart_path),
                electricity_text
            )
        else:
            # Use default image
            self._create_left_text_right_image_slide(
                "Electricity Usage and Energy Conservation Policy",
                electricity_text,
                ENVIRONMENT_IMAGE_MAPPING['ghg_pie']
            )
        
        # Page 15: Energy Efficiency Measures
        efficiency_text = self.content_engine.generate_energy_efficiency_measures(self.config)
        self._create_left_text_right_image_slide(
            "Energy Efficiency Measures",
            efficiency_text,
            ENVIRONMENT_IMAGE_MAPPING['ghg_bar']
        )
        
        print("✓ GHG management pages completed (Pages 13-15)")

    def generate_environmental_management_pages(self):
        """Generate Environmental Management Pages"""
        print("\n[Generating Environmental Management Pages]")
        print("  Pages 16-19: Environmental Management")
        
        # Page 16: 4.6 Green Planting (custom layout: left image 50% wider, right text 40% width)
        plant_text = self.content_engine.generate_green_planting_program(self.config)
        slide = self._add_slide()
        self._add_title(slide, "4.6 Green Planting")
        
        # Left image: CONTENT_WIDTH * 1.5 = 9.0"
        left_image_width = CONTENT_WIDTH * 1.5
        self._add_image(slide, ENVIRONMENT_IMAGE_MAPPING['plant'],
                       left=LEFT_CONTENT_LEFT,
                       top=CONTENT_TOP,
                       width=left_image_width,
                       height=CONTENT_HEIGHT)
        
        # Right text: CONTENT_WIDTH * 0.4 = 2.4"
        right_text_width = CONTENT_WIDTH * 0.4
        right_text_left = LEFT_CONTENT_LEFT + left_image_width + GAP
        self._add_text_box(slide, plant_text,
                          left=right_text_left,
                          top=CONTENT_TOP,
                          width=right_text_width,
                          height=CONTENT_HEIGHT)
        
        # Page 17: 4.7 Water Resource Management
        water_text = self.content_engine.generate_water_management(self.config)
        self._create_left_text_right_image_slide(
            "4.7 Water Resource Management",
            water_text,
            ENVIRONMENT_IMAGE_MAPPING['water']
        )
        
        # Page 18: 4.8 Waste Management
        waste_text = self.content_engine.generate_waste_management(self.config)
        self._create_left_text_right_image_slide(
            "4.8 Waste Management",
            waste_text,
            ENVIRONMENT_IMAGE_MAPPING['waste']
        )
        
        # Page 19: 4.9 Environmental Education and Cooperation (custom layout: left image 50% wider, right text 40% width)
        education_text = self.content_engine.generate_environmental_education(self.config)
        slide = self._add_slide()
        self._add_title(slide, "4.9 Environmental Education and Cooperation")
        
        # Left image: CONTENT_WIDTH * 1.5 = 9.0"
        left_image_width = CONTENT_WIDTH * 1.5
        self._add_image(slide, ENVIRONMENT_IMAGE_MAPPING['ecowork'],
                       left=LEFT_CONTENT_LEFT,
                       top=CONTENT_TOP,
                       width=left_image_width,
                       height=CONTENT_HEIGHT)
        
        # Right text: CONTENT_WIDTH * 0.4 = 2.4"
        right_text_width = CONTENT_WIDTH * 0.4
        right_text_left = LEFT_CONTENT_LEFT + left_image_width + GAP
        self._add_text_box(slide, education_text,
                          left=right_text_left,
                          top=CONTENT_TOP,
                          width=right_text_width,
                          height=CONTENT_HEIGHT)
        
        print("✓ Environmental management pages completed (Pages 16-19)")

    def generate(self):
        """Generate complete Environment Chapter PPTX report"""
        print("\n" + "="*50)
        print("Starting ESG Environment Chapter PPTX Report Generation")
        print(f"Industry: {self.industry}")
        print("="*50)
        
        # Set dynamic emission data
        if self.emission_data:
            try:
                from assets.emission_pptx import set_emission_data
                set_emission_data(self.emission_data)
            except Exception as e:
                print(f"  ⚠ Unable to load dynamic emission data: {e}")
        
        # Generate pages in order:
        # Page 1: Cover
        self.generate_cover_page()
        
        # Pages 2-3: Policy (4.1, 4.2)
        self.generate_policy_pages()
        
        # Pages 5-11: TCFD (7 pages from single PPTX file)
        # Directly insert TCFD PPTX file, no introduction page
        self.generate_tcfd_pages()
        
        # Page 12: SASB Industry Classification (moved after TCFD)
        self.generate_sasb_page()
        
        # Pages 13+: GHG and Environmental Management
        self.generate_ghg_pages()
        self.generate_environmental_management_pages()
        
        print("\n" + "="*50)
        print("PPTX Report Generation Completed!")
        print(f"Total {len(self.prs.slides)} slides")
        print("="*50)
        
        return self.prs

    def save(self, filename):
        """Save PPTX file"""
        self.prs.save(filename)
        print(f"✓ Saved: {filename}")


# For testing
if __name__ == "__main__":
    from datetime import datetime
    
    # Can specify template path
    template_path = r"C:\Users\User\Downloads\templet_resaved.pptx"
    
    engine = EnvironmentPPTXEngine(template_path=template_path)
    pptx = engine.generate()
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"output/ESG_Environment_Chapter_{timestamp}.pptx"
    engine.save(filename)

