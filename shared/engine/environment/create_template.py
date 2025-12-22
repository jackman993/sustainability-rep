"""
建立全新的 A4 橫向模板
包含：右上角半圓形裝飾、右下角浮水印
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# A4 橫向尺寸
SLIDE_WIDTH = Inches(11.69)
SLIDE_HEIGHT = Inches(8.27)

def create_new_template(output_path):
    """建立全新的 A4 橫向模板（在投影片上添加裝飾，不依賴母片）"""
    prs = Presentation()
    
    # 設定 A4 橫向尺寸
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 建立一個測試投影片來添加裝飾（之後可以刪除）
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    # 在投影片上添加裝飾元素（這些會成為模板的一部分）
    # 1. 右上角半圓形（往左 4cm = 1.575"）
    # 計算位置：右邊界 - 半圓半徑 - 4cm
    circle_radius = Inches(1.5)  # 半圓半徑
    circle_center_x = SLIDE_WIDTH - circle_radius - Inches(1.575)  # 往左 4cm
    circle_center_y = Inches(0.5)  # 距離頂部 0.5"
    
    # 建立半圓形（使用橢圓形模擬）
    arc_left = circle_center_x - circle_radius
    arc_top = circle_center_y - circle_radius
    arc_width = circle_radius * 2
    arc_height = circle_radius  # 高度減半，形成半圓效果
    
    # 在投影片上添加橢圓形（模擬半圓）
    arc = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        arc_left,
        arc_top,
        arc_width,
        arc_height
    )
    
    # 設定樣式
    arc.fill.solid()
    arc.fill.fore_color.rgb = RGBColor(230, 240, 230)  # 淡綠色
    arc.line.color.rgb = RGBColor(180, 200, 180)  # 淺綠色邊框
    arc.line.width = Pt(1)
    
    # 2. 右下角浮水印
    watermark_width = Inches(2.5)
    watermark_height = Inches(0.3)
    watermark_left = SLIDE_WIDTH - watermark_width - Inches(0.4)
    watermark_top = SLIDE_HEIGHT - watermark_height - Inches(0.2)
    
    watermark_box = slide.shapes.add_textbox(
        watermark_left, 
        watermark_top, 
        watermark_width, 
        watermark_height
    )
    watermark_frame = watermark_box.text_frame
    watermark_para = watermark_frame.paragraphs[0]
    watermark_para.text = "Sustainability Report"
    watermark_para.font.size = Pt(11)
    watermark_para.font.bold = True
    watermark_para.font.italic = True
    watermark_para.font.name = 'Arial'
    watermark_para.font.color.rgb = RGBColor(200, 200, 200)
    watermark_para.alignment = PP_ALIGN.RIGHT
    
    # 刪除測試投影片（只保留裝飾在母片... 不對，我們需要保留這個投影片作為範例）
    # 實際上，我們建立一個空白模板，裝飾會在環境引擎的 _add_slide() 中添加
    
    # 儲存模板
    prs.save(output_path)
    print(f"✓ 新模板已建立：{output_path}")
    print(f"  - A4 橫向尺寸：{SLIDE_WIDTH.inches}\" x {SLIDE_HEIGHT.inches}\"")
    print(f"  - 右上角半圓形：位置已往左 4cm")
    print(f"  - 右下角浮水印：Sustainability Report")
    
    return output_path


if __name__ == "__main__":
    import sys
    from pathlib import Path
    
    # 輸出路徑
    output_path = Path(__file__).parent / "assets" / "template_new_a4.pptx"
    output_path.parent.mkdir(exist_ok=True)
    
    create_new_template(str(output_path))
    print(f"\n✅ 完成！新模板：{output_path}")

