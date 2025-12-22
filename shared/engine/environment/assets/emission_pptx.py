"""
Emission Engine - PPTX Version
Outputs:
1. Emission Table PPTX
2. Pie Chart Image
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # 使用非互動式後端

# Output Directory
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ============ Default Emission Data ============
DEFAULT_EMISSION_DATA = {
    "data_year": "2024",
    "unit": "tCO₂e",
    "scope1": {
        "gasoline": 1.16,
        "refrigerant": 4.18,
        "subtotal": 5.34
    },
    "scope2": {
        "electricity": 148.11,
        "subtotal": 148.11
    },
    "scope3": {
        "goods_services": 0.00,
        "transportation": 0.00,
        "subtotal": 0.00
    },
    "total": 153.45
}

# Dynamic Data (can be set externally)
EMISSION_DATA = DEFAULT_EMISSION_DATA.copy()

def set_emission_data(data):
    """Set dynamic emission data (passed from Emission Engine)"""
    global EMISSION_DATA
    if data:
        # Get detailed data (prefer gasoline/refrigerant, otherwise split scope1 equally)
        s1_total = data.get("scope1", 0)
        s1_gasoline = data.get("gasoline", s1_total * 0.5)
        s1_refrigerant = data.get("refrigerant", s1_total * 0.5)
        
        s2_total = data.get("scope2", 0)
        s2_electricity = data.get("electricity", s2_total)
        
        total = data.get("total", s1_total + s2_total)
        
        EMISSION_DATA = {
            "data_year": data.get("data_year", "2024"),
            "unit": "tCO₂e",
            "scope1": {
                "gasoline": s1_gasoline,
                "refrigerant": s1_refrigerant,
                "subtotal": s1_total
            },
            "scope2": {
                "electricity": s2_electricity,
                "subtotal": s2_total
            },
            "scope3": {
                "goods_services": 0.00,
                "transportation": 0.00,
                "subtotal": data.get("scope3", 0)
            },
            "total": total
        }
        print(f"  ✓ Dynamic emission data loaded:")
        print(f"    Scope 1: {s1_total:.2f} t (Gasoline {s1_gasoline:.2f}, Refrigerant {s1_refrigerant:.2f})")
        print(f"    Scope 2: {s2_total:.2f} t")
        print(f"    Total Emissions: {total:.2f} tCO₂e")


def set_cell_fill(cell, hex_color):
    """Set cell background color"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = OxmlElement('a:solidFill')
    srgbClr = OxmlElement('a:srgbClr')
    srgbClr.set('val', hex_color)
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def create_emission_table_pptx(output_path=None):
    """Create emission table PPTX (16:9)"""
    prs = Presentation()
    # 16:9 standard presentation size: 13.333" x 7.5"
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Greenhouse Gas Emission Inventory Table (Scope 1-3)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(12), Inches(0.4))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Data Year: {EMISSION_DATA['data_year']} | Unit: {EMISSION_DATA['unit']}"
    p2.font.size = Pt(14)
    p2.font.name = 'Microsoft JhengHei'
    p2.alignment = PP_ALIGN.CENTER
    
    # Create table
    rows = 10
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.4), Inches(10), Inches(5)).table
    
    # Set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.5)
    
    # Header
    headers = ["Scope", "Emission Source", "Emission\n(tCO₂e)", "Notes"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        set_cell_fill(cell, "2F5233")  # Dark green
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows (using dynamic data)
    s1_gas = EMISSION_DATA["scope1"]["gasoline"]
    s1_ref = EMISSION_DATA["scope1"]["refrigerant"]
    s1_sub = EMISSION_DATA["scope1"]["subtotal"]
    s2_elec = EMISSION_DATA["scope2"]["electricity"]
    s2_sub = EMISSION_DATA["scope2"]["subtotal"]
    s3_sub = EMISSION_DATA["scope3"]["subtotal"]
    total = EMISSION_DATA["total"]
    
    # Calculate percentage
    s2_percent = (s2_sub / total * 100) if total > 0 else 0
    
    data_rows = [
        ["Scope 1", "Gasoline (Company Vehicles)", f"{s1_gas:.2f}", "Estimate"],
        ["", "Refrigerant (R-410A)", f"{s1_ref:.2f}", "Maintenance Estimate"],
        ["Subtotal", "", f"{s1_sub:.2f}", ""],
        ["Scope 2", "Purchased Electricity", f"{s2_elec:.2f}", f"{s2_percent:.1f}% of Total"],
        ["Subtotal", "", f"{s2_sub:.2f}", ""],
        ["Scope 3", "Purchased Goods & Services", "0.00", "Not Included"],
        ["", "Transportation", "0.00", "Not Included"],
        ["Subtotal", "", f"{s3_sub:.2f}", ""],
        ["Total", "", f"{total:.2f}", "100%"],
    ]
    
    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = 'Microsoft JhengHei'
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Special row styles
            if row_data[0] in ["Subtotal", "Total"]:
                p.font.bold = True
                if row_data[0] == "Total":
                    set_cell_fill(cell, "E8F5E9")  # Light green
    
    # Save
    if output_path is None:
        output_path = OUTPUT_DIR / "emission_table.pptx"
    
    prs.save(str(output_path))
    print(f"✓ Emission table saved: {output_path}")
    
    return output_path


def create_emission_pie_chart(output_path=None):
    """Create emission pie chart"""
    # Set Chinese font
    plt.rcParams['font.sans-serif'] = ['Microsoft JhengHei', 'SimHei', 'Arial']
    plt.rcParams['axes.unicode_minus'] = False
    
    scope1 = EMISSION_DATA["scope1"]["subtotal"]
    scope2 = EMISSION_DATA["scope2"]["subtotal"]
    scope3 = EMISSION_DATA["scope3"]["subtotal"]
    
    labels = ["Scope 1\n(Direct)", "Scope 2\n(Purchased Electricity)", "Scope 3\n(Other Indirect)"]
    sizes = [scope1, scope2, scope3]
    colors = ["#6BA292", "#007A3D", "#C1C1C1"]
    explode = (0, 0.05, 0)  # Highlight Scope 2
    
    # Avoid issues with 0 values
    if sum(sizes) == 0:
        sizes = [1, 1, 1]
    
    fig, ax = plt.subplots(figsize=(8, 6), dpi=150)
    
    wedges, texts, autotexts = ax.pie(
        sizes,
        labels=labels,
        autopct=lambda p: f"{p:.1f}%\n({p/100*sum([scope1, scope2, scope3]):.2f} t)" if p > 0 else "",
        startangle=90,
        colors=colors,
        explode=explode,
        textprops={"fontsize": 11, "fontweight": "bold"}
    )
    
    ax.set_title("Greenhouse Gas Emission Distribution (tCO₂e)", fontsize=16, fontweight='bold', pad=20)
    ax.axis('equal')
    
    # Legend
    ax.legend(wedges, [f"Scope 1: {scope1:.2f} t", f"Scope 2: {scope2:.2f} t", f"Scope 3: {scope3:.2f} t"],
              title="Emissions", loc="center left", bbox_to_anchor=(1, 0, 0.5, 1))
    
    plt.tight_layout()
    
    if output_path is None:
        output_path = OUTPUT_DIR / "emission_pie_chart.png"
    
    plt.savefig(str(output_path), bbox_inches="tight", dpi=300, facecolor='white')
    plt.close()
    
    print(f"✓ Pie chart saved: {output_path}")
    
    return output_path


def create_emission_table_on_slide_right(slide):
    """Create emission table on right half of slide (reduced 50% width, right-aligned)"""
    # 16:9 width
    slide_w = 13.333
    
    # Original table width 10", reduced 50% = 5"
    table_w = 5.0
    # Right-align: slide_w - table_w - right margin(0.4)
    table_left = slide_w - table_w - 0.4
    
    # Table position: top adjusted to 1.2" to match left text
    rows = 10
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(1.2), Inches(table_w), Inches(5.8))
    tbl = table_shape.table
    
    # Column widths (proportionally reduced 50%)
    tbl.columns[0].width = Inches(1.0)   # Scope
    tbl.columns[1].width = Inches(1.75) # Emission Source
    tbl.columns[2].width = Inches(1.0)   # Emission
    tbl.columns[3].width = Inches(1.25) # Notes
    
    # Header
    headers = ["Scope", "Emission Source", "Emission\n(tCO₂e)", "Notes"]
    for i, header in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = header
        set_cell_fill(cell, "2F5233")  # Dark green
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(9)  # Smaller font
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows (using dynamic data)
    s1_gas = EMISSION_DATA["scope1"]["gasoline"]
    s1_ref = EMISSION_DATA["scope1"]["refrigerant"]
    s1_sub = EMISSION_DATA["scope1"]["subtotal"]
    s2_elec = EMISSION_DATA["scope2"]["electricity"]
    s2_sub = EMISSION_DATA["scope2"]["subtotal"]
    s3_sub = EMISSION_DATA["scope3"]["subtotal"]
    total = EMISSION_DATA["total"]
    
    # Calculate percentage
    s2_percent = (s2_sub / total * 100) if total > 0 else 0
    
    data_rows = [
        ["Scope 1", "Gasoline (Company Vehicles)", f"{s1_gas:.2f}", "Estimate"],
        ["", "Refrigerant (R-410A)", f"{s1_ref:.2f}", "Maintenance Estimate"],
        ["Subtotal", "", f"{s1_sub:.2f}", ""],
        ["Scope 2", "Purchased Electricity", f"{s2_elec:.2f}", f"{s2_percent:.1f}% of Total"],
        ["Subtotal", "", f"{s2_sub:.2f}", ""],
        ["Scope 3", "Purchased Goods & Services", "0.00", "Not Included"],
        ["", "Transportation", "0.00", "Not Included"],
        ["Subtotal", "", f"{s3_sub:.2f}", ""],
        ["Total", "", f"{total:.2f}", "100%"],
    ]
    
    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, value in enumerate(row_data):
            cell = tbl.cell(row_idx, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)  # Smaller font
            p.font.name = 'Microsoft JhengHei'
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Special row styles
            if row_data[0] in ["Subtotal", "Total"]:
                p.font.bold = True
                if row_data[0] == "Total":
                    set_cell_fill(cell, "E8F5E9")  # Light green
    
    print("  ✓ Emission table inserted (right half, reduced 50%)")


def generate_all():
    """Generate all emission-related files"""
    print("\n[Generating Emission Engine Output]")
    table_path = create_emission_table_pptx()
    chart_path = create_emission_pie_chart()
    
    return {
        "table_pptx": table_path,
        "pie_chart": chart_path
    }


if __name__ == "__main__":
    results = generate_all()
    print(f"\nCompleted!")
    print(f"  Table: {results['table_pptx']}")
    print(f"  Chart: {results['pie_chart']}")

