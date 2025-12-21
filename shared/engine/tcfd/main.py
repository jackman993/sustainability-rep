"""
TCFD Main Engine
主邏輯：協調配置、內容生成和表格生成
"""
import os
import sys
import importlib.util
from pathlib import Path
from typing import Dict, List, Optional, Any
import json
import re

from . import config
from . import content
from ..path_manager import get_tcfd_output_path, update_session_activity

# 嘗試導入 Claude API
try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False


def load_table_module(script_file: str):
    """
    動態載入表格生成模組
    
    Args:
        script_file: 表格腳本文件名（如 'table01.py'）
    
    Returns:
        載入的模組對象
    """
    script_path = config.TABLES_DIR / script_file
    
    if not script_path.exists():
        raise FileNotFoundError(f"Table script not found: {script_path}")
    
    module_name = f"tcfd_table_{script_file.replace('.py', '')}"
    spec = importlib.util.spec_from_file_location(module_name, script_path)
    
    if spec is None:
        raise ImportError(f"Cannot load module: {script_path}")
    
    module = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = module
    spec.loader.exec_module(module)
    
    return module


def call_claude_api(prompt: str, api_key: str, model: str = None) -> str:
    """
    調用 Claude API 生成內容
    
    使用多個備選模型，第一個失敗就嘗試下一個
    
    Args:
        prompt: 完整的 prompt
        api_key: Anthropic API Key
        model: Claude 模型名稱（可選，如果不提供則使用備選列表）
    
    Returns:
        API 返回的文本內容
    """
    if not ANTHROPIC_AVAILABLE:
        raise ImportError("anthropic package not installed. Install with: pip install anthropic")
    
    # 備選模型列表（按優先順序）
    model_list = [
        "claude-3-5-sonnet-20240620",
        "claude-3-opus-20240229",
        "claude-3-sonnet-20240229",
        "claude-3-haiku-20240307"
    ]
    
    # 如果指定了模型，優先使用
    if model:
        model_list = [model] + [m for m in model_list if m != model]
    
    client = anthropic.Anthropic(api_key=api_key)
    
    # 嘗試每個模型，直到成功
    last_error = None
    for model_name in model_list:
        try:
            message = client.messages.create(
                model=model_name,
                max_tokens=2000,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            return message.content[0].text
        except Exception as e:
            last_error = e
            # 如果是模型不存在的錯誤，嘗試下一個模型
            if "not_found_error" in str(e) or "404" in str(e):
                continue
            # 其他錯誤直接拋出
            raise
    
    # 所有模型都失敗
    raise Exception(f"All models failed. Last error: {str(last_error)}")


def parse_llm_response(response: str) -> List[str]:
    """
    解析 LLM 返回的內容，提取表格行數據（優化版本，強制 ||| 格式）
    
    Args:
        response: LLM 返回的文本
    
    Returns:
        表格內容列表（每行是 ||| 分隔的字符串）
    """
    lines = []
    # 尋找包含 ||| 分隔符的行（優先）
    for line in response.split('\n'):
        line = line.strip()
        if '|||' in line:
            # 清理可能的編號前綴（如 "Line 1:", "1.", 等）
            cleaned = re.sub(r'^(Line\s*\d+[:.]?\s*|^\d+[.:]\s*)', '', line, flags=re.IGNORECASE)
            if cleaned.strip():
                lines.append(cleaned.strip())
    
    # 如果沒有找到 ||| 分隔的行，嘗試其他格式
    if not lines:
        # 嘗試尋找列表格式
        for line in response.split('\n'):
            line = line.strip()
            if line and (line.startswith('-') or line.startswith('•') or re.match(r'^\d+\.', line)):
                # 嘗試解析為表格格式
                cleaned = re.sub(r'^[-•0-9. ]+', '', line)
                if cleaned:
                    lines.append(cleaned)
    
    # 如果還是沒有，返回原始響應的第一部分（添加 ||| 分隔符）
    if not lines:
        # 嘗試從原始響應中提取前200字符，並添加分隔符
        first_part = response[:200].strip()
        if first_part:
            lines.append(first_part + " ||| N/A ||| N/A")
        else:
            lines.append("N/A ||| N/A ||| N/A")
    
    # 返回最多 10 行（通常只需要 2 行）
    return lines[:10]


def generate_mock_data(prompt_id: str, industry: str = None, carbon_emission: Dict[str, Any] = None) -> List[str]:
    """
    生成模擬數據（用於測試）
    
    Args:
        prompt_id: Prompt ID
        industry: 產業名稱
        carbon_emission: 碳排放數據
    
    Returns:
        模擬的表格內容列表
    """
    industry = industry or "Manufacturing"
    
    # 根據不同的表格類型生成不同的模擬數據（增加行數和文字密度）
    if 'trans' in prompt_id:
        return [
            # Line 1: Policy & Regulation Risk (對應 Row 2) - 詳細版本
            f"Policy & Regulation Risk;{industry} sector faces comprehensive regulatory changes including new carbon tax regulations, stricter environmental compliance requirements, and sector-specific emission standards. Government policies mandate 30% emission reduction by 2030, requiring immediate action. Additional regulations include product certification requirements, supply chain transparency mandates, and environmental permit restrictions. Compliance deadlines are approaching with enforcement mechanisms including fines, operational restrictions, and market access limitations ||| Financial Impact: $500K-1M annually in carbon taxes, plus $200K-400K in compliance costs including monitoring systems, reporting, and certification fees. Potential penalties up to $500K for non-compliance. Additional costs include $300K-600K for process modifications and $150K-300K for staff training. Total first-year compliance cost estimated at $1.5M-2.5M ||| Mitigation: Invest in carbon offset programs ($200K initial) and energy efficiency upgrades ($500K). Implement comprehensive compliance monitoring system ($300K). Establish regulatory tracking and reporting framework ($150K). Develop compliance training programs ($100K). Budget: $1.25M initial investment, $350K annually for maintenance, monitoring, and ongoing compliance",
            # Line 2: Green Product & Technology Risk (對應 Row 3) - 詳細版本
            f"Green Product & Technology Risk;Significant technology disruption in {industry} sector with shift to green technologies, sustainable materials, and low-carbon production processes. Competitors adopting advanced green technologies including renewable energy systems, circular economy solutions, and digital automation gaining competitive advantage. Emerging technologies such as AI-driven efficiency optimization, blockchain for supply chain transparency, and advanced recycling processes transforming {industry} operations ||| Financial Impact: Technology obsolescence risk $1M-2M for existing equipment. Critical need for R&D investment $1.5M-3M to stay competitive. Market share erosion 10-15% if no action, representing $2M-4M revenue risk. Additional costs include $800K-1.2M for new equipment, $400K-600K for technology integration, and $200K-400K for staff retraining. ROI from technology adoption estimated at 15-25% over 5 years ||| Mitigation: Accelerate green technology R&D program ($1.5M over 3 years). Develop sustainable product lines ($800K). Establish strategic partnerships with green tech suppliers ($200K). Implement technology pilot programs ($300K). Create innovation lab for testing new technologies ($400K). Budget: $3.2M total investment over 3 years, $600K annually for R&D and technology maintenance"
        ]
    elif 'physical' in prompt_id:
        return [
            f"Acute Risk (Short Term);Extreme weather events including floods, storms, and heatwaves affecting {industry} operations. Increased frequency of natural disasters disrupts supply chains and damages facilities ||| Financial Impact: $300K-800K per event in direct damages. Business interruption costs of $200K-500K per event. Insurance claims may cover only 60-70% of losses ||| Mitigation: Strengthen facility resilience with upgraded infrastructure. Implement emergency response protocols and business continuity plans. Budget: $500K for infrastructure upgrades, $100K annually for maintenance",
            f"Chronic Risk (Long Term);Rising temperatures and changing precipitation patterns impacting {industry} supply chain. Water scarcity and heat stress affecting operations and agricultural inputs ||| Financial Impact: $1M-3M annually in increased operational costs. Supply chain disruptions costing $500K-1.5M per year. Reduced productivity due to heat stress ||| Mitigation: Diversify supply sources geographically and invest in climate-resilient infrastructure. Develop water conservation and heat management systems. Budget: $800K initial investment, $300K annually for operations",
            f"Sea Level Rise Risk;Coastal facilities and transportation infrastructure at risk from rising sea levels. Port operations and logistics networks vulnerable to flooding and storm surges ||| Financial Impact: $2M-5M in potential relocation costs. Infrastructure damage of $500K-1.5M over 10 years. Increased insurance premiums of $200K-400K annually ||| Mitigation: Conduct comprehensive risk assessment and develop relocation or protection strategies. Invest in flood protection infrastructure. Budget: $1.2M for assessment and planning, $600K for protection measures",
            f"Biodiversity & Ecosystem Risk;Loss of ecosystem services affecting {industry} operations. Reduced water quality, soil degradation, and loss of natural resources impacting long-term sustainability ||| Financial Impact: $400K-1M annually in increased resource costs. Need for alternative resource sourcing costing $300K-600K. Potential regulatory restrictions on resource use ||| Mitigation: Implement biodiversity conservation programs and sustainable resource management practices. Engage in ecosystem restoration initiatives. Budget: $500K for conservation programs, $200K annually for maintenance"
        ]
    elif 'opp_resource' in prompt_id:
        return [
            f"Resource Efficiency;Implement energy-efficient processes and waste reduction systems in {industry} operations. Upgrade to high-efficiency equipment and optimize production workflows to minimize resource consumption ||| Financial Benefit: $400K-600K annual savings in energy costs. Reduced material waste saving $200K-400K annually. Lower operational expenses improving profit margins by 5-8% ||| Investment: Upgrade equipment and implement process optimization systems. Install energy monitoring and management systems. Budget: $1.5M initial investment, payback period 2-3 years",
            f"Energy Source;Transition to renewable energy sources including solar and wind power for {industry} facilities. Reduce dependence on fossil fuels and achieve carbon neutrality goals ||| Financial Benefit: $200K-400K annual savings in energy costs. Long-term price stability reducing energy cost volatility. Potential carbon credit revenue of $50K-150K annually ||| Investment: Solar panel installation and wind power infrastructure. Energy storage systems and grid integration. Budget: $2M initial investment, payback period 5-7 years",
            f"Water Management;Implement advanced water recycling and conservation systems. Reduce water consumption through process optimization and reuse technologies ||| Financial Benefit: $150K-300K annual savings in water costs. Reduced wastewater treatment costs saving $100K-200K. Lower regulatory compliance costs ||| Investment: Water recycling systems and process optimization. Monitoring and control systems. Budget: $800K initial investment, payback period 3-4 years",
            f"Circular Economy;Develop closed-loop production systems and material recycling programs. Minimize waste generation and maximize resource recovery from production processes ||| Financial Benefit: $300K-500K annual savings from reduced waste disposal costs. Revenue from recycled materials of $100K-200K annually. Reduced raw material costs ||| Investment: Recycling infrastructure and process redesign. Material recovery and sorting systems. Budget: $1.2M initial investment, payback period 2-3 years"
        ]
    elif 'opp_product' in prompt_id:
        return [
            f"Products & Services;Develop sustainable {industry} products for growing green market. Create eco-friendly alternatives that meet consumer demand for environmentally responsible options ||| Financial Benefit: $5M-10M new revenue from green product lines over 3 years. Market share growth of 10-20% in sustainable segment. Premium pricing potential of 15-25% ||| Investment: R&D for sustainable product development and marketing campaigns. Production line modifications and certification processes. Budget: $3M initial investment, $1M annually for R&D and marketing",
            f"New Markets & Resilience;Expand {industry} business to climate-resilient regions and emerging markets. Diversify geographic presence to reduce climate risk exposure ||| Financial Benefit: $3M-7M new revenue from market expansion over 5 years. Reduced climate risk through geographic diversification. Access to new customer segments ||| Investment: Market entry and setup costs. Local partnerships and distribution networks. Budget: $2M initial investment, $500K annually for operations",
            f"Green Technology Innovation;Invest in research and development of climate-friendly technologies. Create proprietary solutions that address environmental challenges in {industry} sector ||| Financial Benefit: $2M-5M potential revenue from technology licensing. Competitive advantage in green technology space. Potential for industry leadership position ||| Investment: R&D facilities and research teams. Technology development and patent applications. Budget: $2.5M initial investment, $1.2M annually for R&D",
            f"Sustainable Supply Chain Services;Develop consulting and advisory services for supply chain sustainability. Leverage expertise to help other companies improve their environmental performance ||| Financial Benefit: $1M-3M annual revenue from consulting services. Recurring revenue model with high margins. Enhanced reputation and brand value ||| Investment: Service development and team building. Marketing and business development activities. Budget: $800K initial investment, $400K annually for operations"
        ]
    elif 'metrics' in prompt_id:
        total_emission = carbon_emission.get('total_tco2e', 100) if carbon_emission else 100
        return [
            f"GHG Emissions Target;Reduce total emissions by 30% by 2030 (Current: {total_emission:.1f} tCO2e). Achieve carbon neutrality by 2040 through comprehensive emission reduction strategies ||| Progress: 15% reduction achieved through energy efficiency improvements. On track to meet 2030 target with current initiatives. Additional 10% reduction planned through renewable energy transition ||| Action Plan: Energy efficiency projects, renewable energy adoption, and process optimization. Implement carbon capture technologies where feasible. Budget: $1M for efficiency projects, $2M for renewable energy, $500K for monitoring",
            f"Renewable Energy Target;Increase renewable energy usage to 50% by 2028 and 100% by 2035. Reduce dependence on fossil fuels and achieve energy independence ||| Progress: 25% renewable energy achieved through solar installations. Wind power projects in planning phase. Grid integration systems being developed ||| Action Plan: Solar and wind investments, energy storage systems, and grid infrastructure upgrades. Develop long-term energy procurement strategy. Budget: $2M for solar, $1.5M for wind, $800K for storage systems",
            f"Water Conservation Target;Reduce water consumption by 40% by 2030 through efficiency improvements and recycling systems. Achieve zero wastewater discharge by 2035 ||| Progress: 20% reduction achieved through process optimization. Water recycling pilot program showing positive results. Monitoring systems in place ||| Action Plan: Advanced water treatment and recycling systems. Process redesign for water efficiency. Implement water monitoring and management systems. Budget: $1.2M for recycling systems, $600K for process upgrades, $200K for monitoring",
            f"Waste Reduction Target;Achieve zero waste to landfill by 2030 through comprehensive recycling and circular economy initiatives. Reduce overall waste generation by 50% ||| Progress: 30% waste reduction achieved through source reduction programs. Recycling rate increased to 75%. Composting programs established ||| Action Plan: Waste minimization programs, advanced recycling infrastructure, and circular economy partnerships. Develop waste-to-energy solutions. Budget: $900K for recycling infrastructure, $400K for programs, $300K for partnerships"
        ]
    elif 'systemic' in prompt_id:
        return [
            f"Industry Certification;Obtain ISO 14001 environmental management certification and other relevant industry standards. Demonstrate commitment to environmental excellence and regulatory compliance ||| Driver: Customer requirements and market expectations. Competitive advantage in B2B markets. Regulatory compliance needs ||| Gap: Missing documentation and formalized environmental management systems. Lack of certified processes and procedures. Need for staff training ||| Action: Complete certification process with comprehensive documentation. Implement environmental management systems. Train staff on requirements. Budget: $150K for certification, $50K annually for maintenance",
            f"Supply Chain Visibility;Enhance visibility into {industry} supply chain emissions and environmental impacts. Develop comprehensive supplier assessment and monitoring programs ||| Driver: Regulatory compliance and stakeholder expectations. Risk management and brand protection. Customer transparency requirements ||| Gap: Limited supplier data and incomplete emissions tracking. Lack of standardized reporting from suppliers. Insufficient monitoring systems ||| Action: Implement supplier reporting system and assessment framework. Develop data collection and analysis capabilities. Establish supplier engagement programs. Budget: $300K for systems, $100K annually for operations",
            f"Climate Risk Disclosure;Develop comprehensive climate risk disclosure framework aligned with TCFD recommendations. Enhance transparency and stakeholder communication ||| Driver: Regulatory requirements and investor expectations. TCFD framework adoption. ESG rating improvements ||| Gap: Incomplete risk assessment and limited disclosure capabilities. Lack of standardized reporting formats. Insufficient data collection systems ||| Action: Develop risk assessment methodologies and reporting frameworks. Implement data collection and analysis systems. Train staff on disclosure requirements. Budget: $250K for framework development, $80K annually for reporting",
            f"Stakeholder Engagement;Establish comprehensive stakeholder engagement program for climate-related issues. Build relationships with investors, customers, and communities ||| Driver: Reputation management and risk mitigation. Regulatory compliance. Market positioning and competitive advantage ||| Gap: Limited stakeholder communication channels. Lack of formal engagement processes. Insufficient climate communication capabilities ||| Action: Develop stakeholder engagement strategy and communication plans. Establish regular dialogue mechanisms. Create climate communication materials. Budget: $180K for program development, $60K annually for activities"
        ]
    elif 'resilience' in prompt_id:
        return [
            f"Workforce Capability;Train {industry} workforce on climate adaptation strategies and resilience planning. Develop skills for managing climate-related disruptions and operational challenges ||| Stressor: Climate-related disruptions affecting operations and supply chains. Extreme weather events causing operational delays. Changing regulatory and market conditions ||| Impact: Operational delays and productivity losses. Increased costs from disruption management. Need for adaptive workforce capabilities ||| Action: Comprehensive training program covering climate risks, adaptation strategies, and resilience planning. Develop internal expertise and knowledge base. Budget: $200K for training programs, $80K annually for ongoing development",
            f"Supply Chain Security;Diversify {industry} supply sources to reduce climate risks and single-point-of-failure vulnerabilities. Build resilient supply network with geographic and supplier diversification ||| Stressor: Single-source dependency creating vulnerability. Climate risks affecting key suppliers. Supply disruptions from extreme weather events ||| Impact: Supply disruptions causing production delays. Increased costs from emergency sourcing. Risk of customer service failures ||| Action: Identify and onboard alternative suppliers in different geographic regions. Develop supply chain risk assessment and monitoring systems. Budget: $500K for supplier development, $150K annually for monitoring",
            f"Operational Flexibility;Develop flexible operational systems that can adapt to climate-related disruptions. Implement agile production and logistics capabilities ||| Stressor: Climate variability affecting production schedules. Extreme weather disrupting operations. Changing resource availability ||| Impact: Production inefficiencies and increased costs. Difficulty maintaining service levels during disruptions. Operational constraints ||| Action: Implement flexible production systems and agile logistics. Develop contingency planning and rapid response capabilities. Budget: $600K for system development, $200K annually for maintenance",
            f"Financial Resilience;Build financial reserves and risk management capabilities to withstand climate-related financial impacts. Develop insurance and risk transfer strategies ||| Stressor: Climate-related financial losses from disruptions. Increased insurance costs. Market volatility affecting financial stability ||| Impact: Financial losses from climate events. Increased costs reducing profitability. Need for financial buffer and risk management ||| Action: Develop financial risk management strategies and insurance programs. Build financial reserves for climate contingencies. Budget: $400K for risk management, $150K annually for insurance and reserves"
        ]
    else:
        # 默認模擬數據
        return [
            f"{industry} Item A;Detail 1 ||| Impact $100K ||| Action Plan A;Budget $50K",
            f"{industry} Item B;Detail 2 ||| Impact $200K ||| Action Plan B;Budget $80K"
        ]


def generate_table_content(
    prompt_id: str,
    industry: str = None,
    revenue: str = None,
    carbon_emission: Dict[str, Any] = None,
    llm_api_key: str = None,
    llm_provider: str = None,
    use_mock: bool = False
) -> List[str]:
    """
    生成表格內容（調用 LLM 或使用模擬數據）
    
    Args:
        prompt_id: Prompt ID（如 'prompt_table_1_trans'）
        industry: 產業名稱
        revenue: 營收
        carbon_emission: 碳排放數據字典
        llm_api_key: LLM API Key（可選）
        llm_provider: LLM 提供商（可選，如 'anthropic'）
        use_mock: 是否使用模擬數據（True = Mock, False = 嘗試 API）
    
    Returns:
        表格內容列表（每行是 ||| 分隔的字符串）
    """
    # 構建完整 prompt
    full_prompt = content.get_prompt(
        prompt_id=prompt_id,
        industry=industry,
        revenue=revenue,
        carbon_emission=carbon_emission
    )
    
    # 如果明確要求使用 Mock，或沒有提供 API Key，使用模擬數據
    if use_mock or not llm_api_key or not llm_provider:
        return generate_mock_data(prompt_id, industry, carbon_emission)
    
    # 嘗試調用 LLM API
    try:
        if llm_provider.lower() == 'anthropic' or llm_provider.lower() == 'claude':
            response = call_claude_api(full_prompt, llm_api_key)
            return parse_llm_response(response)
        else:
            # 不支持的提供商，使用 Mock
            return generate_mock_data(prompt_id, industry, carbon_emission)
    except Exception as e:
        print(f"Error calling LLM API: {str(e)}")
        # API 調用失敗，回退到 Mock
        return generate_mock_data(prompt_id, industry, carbon_emission)


def generate_table(
    page_key: str,
    output_dir: Path = None,
    industry: str = None,
    revenue: str = None,
    carbon_emission: Dict[str, Any] = None,
    llm_api_key: str = None,
    llm_provider: str = None,
    use_mock: bool = False
) -> Optional[Path]:
    """
    生成單個 TCFD 表格
    
    Args:
        page_key: 頁面鍵（如 'page_1'）
        output_dir: 輸出目錄（默認為 config.OUTPUT_DIR）
        industry: 產業名稱
        revenue: 營收
        carbon_emission: 碳排放數據
        llm_api_key: LLM API Key
        llm_provider: LLM 提供商
    
    Returns:
        生成的 PowerPoint 文件路徑，失敗返回 None
    """
    if page_key not in config.TCFD_PAGES:
        raise ValueError(f"Invalid page_key: {page_key}")
    
    page_info = config.TCFD_PAGES[page_key]
    output_dir = output_dir or config.OUTPUT_DIR
    output_dir.mkdir(exist_ok=True)
    
    try:
        # 1. 載入表格生成模組
        table_module = load_table_module(page_info['script_file'])
        
        # 2. 生成表格內容
        # 如果沒有明確指定 use_mock，則根據是否有 API key 決定
        if use_mock is None:
            use_mock = (not llm_api_key or not llm_provider)
        
        data_lines = generate_table_content(
            prompt_id=page_info['prompt_id'],
            industry=industry,
            revenue=revenue,
            carbon_emission=carbon_emission,
            llm_api_key=llm_api_key,
            llm_provider=llm_provider,
            use_mock=use_mock
        )
        
        # 3. 調用表格生成函數
        entry_func_name = page_info['entry_function']
        
        if not hasattr(table_module, entry_func_name):
            raise AttributeError(
                f"Function '{entry_func_name}' not found in {page_info['script_file']}"
            )
        
        func = getattr(table_module, entry_func_name)
        
        # 4. 設定輸出檔名
        safe_script_name = page_info['script_file'].replace('.py', '')
        out_filename = f"TCFD_{page_key}_{safe_script_name}.pptx"
        out_path = output_dir / out_filename
        
        # 5. 生成表格（傳入數據和完整路徑）
        func(data_lines, filename=str(out_path))
        
        return out_path
        
    except Exception as e:
        print(f"Error generating table {page_key}: {str(e)}")
        return None


def generate_all_tables(
    output_dir: Path = None,
    industry: str = None,
    revenue: str = None,
    carbon_emission: Dict[str, Any] = None,
    llm_api_key: str = None,
    llm_provider: str = None
) -> Dict[str, Optional[Path]]:
    """
    生成所有 TCFD 表格
    
    Args:
        output_dir: 輸出目錄
        industry: 產業名稱
        revenue: 營收
        carbon_emission: 碳排放數據
        llm_api_key: LLM API Key
        llm_provider: LLM 提供商
    
    Returns:
        字典：{page_key: 文件路徑}
    """
    results = {}
    
    for page_key in config.TCFD_PAGES.keys():
        results[page_key] = generate_table(
            page_key=page_key,
            output_dir=output_dir,
            industry=industry,
            revenue=revenue,
            carbon_emission=carbon_emission,
            llm_api_key=llm_api_key,
            llm_provider=llm_provider
        )
    
    return results


def generate_combined_pptx(
    output_filename: str = "TCFD_table.pptx",
    template_path: Path = None,
    industry: str = None,
    revenue: str = None,
    carbon_emission: Dict[str, Any] = None,
    llm_api_key: str = None,
    llm_provider: str = None,
    use_mock: bool = False
) -> Optional[Path]:
    """
    生成包含所有表格的單個 PPTX 文件
    
    統一處理邏輯：
    - table0607 (page_6, page_7): 接受 prs 參數，直接添加到主 PPTX
    - table01-05 (page_1 到 page_5): 創建新文件，然後複製 slide
    
    Args:
        output_filename: 輸出文件名（已棄用，現在使用 path_manager 的標準文件名）
        template_path: 模板文件路徑（handdrawppt.pptx）
        industry: 產業名稱
        revenue: 營收
        carbon_emission: 碳排放數據
        llm_api_key: LLM API Key
        llm_provider: LLM 提供商
        use_mock: 是否使用模擬數據
    
    Returns:
        生成的 PowerPoint 文件路徑（位於 output/{session_id}/TCFD_table.pptx）
        
    Note:
        輸出路徑現在通過 path_manager.get_tcfd_output_path() 統一管理，
        自動包含 session_id，並會更新會話活動時間。
    """
    try:
        from pptx import Presentation
        import tempfile
        import os
        
        # 載入模板
        if template_path and template_path.exists():
            prs = Presentation(str(template_path))
        else:
            # 使用默認模板路徑
            default_template = config.BASE_DIR / "handdrawppt.pptx"
            if default_template.exists():
                prs = Presentation(str(default_template))
            else:
                prs = Presentation()
        
        # 刪除模板中的所有預設 slide（避免空白頁）
        # 因為我們要一次輸出7頁，不需要模板的預設頁面
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # 按順序生成每個表格
        slide_order = ['page_1', 'page_2', 'page_3', 'page_4', 'page_5', 'page_6', 'page_7']
        
        for page_key in slide_order:
            if page_key not in config.TCFD_PAGES:
                continue
            
            page_info = config.TCFD_PAGES[page_key]
            
            # 載入表格生成模組
            table_module = load_table_module(page_info['script_file'])
            func = getattr(table_module, page_info['entry_function'])
            
            # 檢查函數簽名
            import inspect
            sig = inspect.signature(func)
            has_prs_param = 'prs' in sig.parameters
            
            # 生成表格內容（所有表格都需要）
            data_lines = generate_table_content(
                prompt_id=page_info['prompt_id'],
                industry=industry,
                revenue=revenue,
                carbon_emission=carbon_emission,
                llm_api_key=llm_api_key,
                llm_provider=llm_provider,
                use_mock=use_mock
            )
            
            # 統一處理：如果函數接受 prs 參數，直接傳入；否則使用臨時文件
            try:
                if has_prs_param:
                    # 函數接受 prs 參數，直接添加到主 PPTX（一次輸出，不用 add_slide）
                    # 檢查是否還需要 data_lines 參數
                    if 'data_lines' in sig.parameters:
                        func(prs=prs, data_lines=data_lines)
                    else:
                        func(prs=prs)
                else:
                    # 函數不接受 prs 參數，使用臨時文件（向後兼容）
                    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                        temp_file = Path(tmp.name)
                    
                    try:
                        # 生成臨時文件
                        func(data_lines=data_lines, filename=str(temp_file))
                        
                        # 讀取臨時文件的 slide，直接移動到主 PPTX
                        temp_prs = Presentation(str(temp_file))
                        
                        # 直接移動 slide 的 XML 元素到主 prs（一次輸出，不用 add_slide）
                        import copy
                        for slide in temp_prs.slides:
                            # 深拷貝 slide 的 XML 元素
                            slide_element = copy.deepcopy(slide._element)
                            # 添加到主 prs 的 slide ID 列表
                            prs.slides._sldIdLst.append(slide_element)
                    finally:
                        # 刪除臨時文件
                        if temp_file.exists():
                            os.unlink(temp_file)
            except Exception as table_error:
                # 捕獲單個表格的錯誤，提供詳細信息
                error_msg = f"Error generating table {page_key} ({page_info['title']}): {str(table_error)}"
                print(error_msg)
                import traceback
                print(f"Table {page_key} traceback:")
                traceback.print_exc()
                # 重新拋出錯誤，讓外層處理
                raise Exception(error_msg) from table_error
        
        # 保存文件 - 使用統一的路徑管理器
        import streamlit as st
        
        print("[DEBUG] ========== TCFD 文件保存開始 ==========")
        
        # 使用統一的 path_manager 獲取輸出路徑（包含 session_id）
        try:
            print("[DEBUG] Step 1: 獲取輸出路徑...")
            output_path = get_tcfd_output_path()
            print(f"[DEBUG] Step 1 成功: {output_path}")
        except Exception as path_error:
            error_msg = f"[ERROR] Failed to get output path: {str(path_error)}"
            print(error_msg)
            import traceback
            full_traceback = traceback.format_exc()
            print(full_traceback)
            # 嘗試在 Streamlit 中顯示錯誤（如果可用）
            try:
                st.error(f"❌ 路徑獲取失敗: {str(path_error)}")
                with st.expander("詳細錯誤信息", expanded=True):
                    st.code(full_traceback)
            except:
                pass
            raise Exception(error_msg) from path_error
        
        # 確保目錄存在（強制創建）
        try:
            print(f"[DEBUG] Step 2: 強制創建目錄 {output_path.parent}...")
            
            # 方法 1: 標準 mkdir
            try:
                output_path.parent.mkdir(parents=True, exist_ok=True)
            except Exception as e1:
                print(f"[WARNING] mkdir() 失敗: {e1}")
                # 方法 2: 使用 os.makedirs 強制創建
                try:
                    os.makedirs(str(output_path.parent), exist_ok=True, mode=0o777)
                    print(f"[DEBUG] os.makedirs() 成功")
                except Exception as e2:
                    print(f"[ERROR] os.makedirs() 也失敗: {e2}")
                    raise
            
            # 驗證目錄是否真的存在
            if not output_path.parent.exists():
                raise Exception(f"目錄創建失敗，目錄不存在: {output_path.parent}")
            
            # 嘗試修改權限確保可寫
            try:
                os.chmod(str(output_path.parent), 0o777)
            except:
                pass  # 權限修改失敗不影響
            
            print(f"[DEBUG] Step 2 成功: 目錄已創建/存在")
            print(f"[DEBUG] 目錄是否存在: {output_path.parent.exists()}")
            print(f"[DEBUG] 目錄是否可寫: {os.access(output_path.parent, os.W_OK)}")
        except Exception as dir_error:
            error_msg = f"[ERROR] Failed to create output directory: {str(dir_error)}"
            print(error_msg)
            import traceback
            full_traceback = traceback.format_exc()
            print(full_traceback)
            try:
                st.error(f"❌ 目錄創建失敗: {str(dir_error)}")
                with st.expander("詳細錯誤信息", expanded=True):
                    st.code(full_traceback)
            except:
                pass
            raise Exception(error_msg) from dir_error
        
        # 保存文件
        try:
            print(f"[DEBUG] Step 3: 保存文件到 {output_path}...")
            print(f"[DEBUG] Presentation 對象: {type(prs)}")
            print(f"[DEBUG] Slides 數量: {len(prs.slides)}")
            
            # 確保路徑是字符串
            save_path = str(output_path)
            print(f"[DEBUG] 保存路徑 (字符串): {save_path}")
            
            # 強制保存文件（多次嘗試）
            print(f"[DEBUG] 準備保存文件到: {save_path}")
            
            # 確保父目錄存在且可寫（多次嘗試）
            parent_dir = Path(save_path).parent
            if not parent_dir.exists():
                print(f"[WARNING] 父目錄不存在，強制創建: {parent_dir}")
                # 方法 1: Path.mkdir
                try:
                    parent_dir.mkdir(parents=True, exist_ok=True)
                except Exception as e1:
                    print(f"[WARNING] Path.mkdir 失敗: {e1}")
                    # 方法 2: os.makedirs
                    try:
                        os.makedirs(str(parent_dir), exist_ok=True, mode=0o777)
                    except Exception as e2:
                        print(f"[ERROR] os.makedirs 也失敗: {e2}")
                        raise Exception(f"無法創建父目錄: {e1}, {e2}")
            
            # 驗證父目錄存在
            if not parent_dir.exists():
                raise Exception(f"父目錄創建失敗，目錄不存在: {parent_dir}")
            
            # 檢查並修復權限
            if not os.access(parent_dir, os.W_OK):
                print(f"[WARNING] 父目錄不可寫，嘗試修改權限...")
                try:
                    os.chmod(str(parent_dir), 0o777)
                    print(f"[DEBUG] 權限修改完成")
                except Exception as perm_ex:
                    print(f"[WARNING] 無法修改權限: {perm_ex}")
                    # 即使權限修改失敗，也繼續嘗試保存
            
            print(f"[DEBUG] 父目錄狀態: 存在={parent_dir.exists()}, 可寫={os.access(parent_dir, os.W_OK)}")
            
            # 強制保存文件（多次嘗試不同方法）
            save_success = False
            last_error = None
            
            # 嘗試 1: 直接保存
            try:
                print(f"[DEBUG] 嘗試 1: 直接保存到 {save_path}")
                prs.save(save_path)
                print(f"[DEBUG] prs.save() 調用完成")
                save_success = True
            except Exception as save_ex1:
                print(f"[WARNING] 嘗試 1 失敗: {save_ex1}")
                last_error = save_ex1
                
                # 嘗試 2: 使用絕對路徑
                try:
                    abs_path = Path(save_path).resolve()
                    print(f"[DEBUG] 嘗試 2: 使用絕對路徑 {abs_path}")
                    prs.save(str(abs_path))
                    save_path = str(abs_path)  # 更新路徑
                    print(f"[DEBUG] 使用絕對路徑保存成功")
                    save_success = True
                except Exception as save_ex2:
                    print(f"[WARNING] 嘗試 2 也失敗: {save_ex2}")
                    last_error = save_ex2
                    
                    # 嘗試 3: 先刪除舊文件（如果存在）再保存
                    try:
                        if Path(save_path).exists():
                            print(f"[DEBUG] 嘗試 3: 刪除舊文件後保存")
                            Path(save_path).unlink()
                        prs.save(save_path)
                        print(f"[DEBUG] 刪除舊文件後保存成功")
                        save_success = True
                    except Exception as save_ex3:
                        print(f"[ERROR] 嘗試 3 也失敗: {save_ex3}")
                        last_error = save_ex3
            
            if not save_success:
                raise Exception(f"所有保存方法都失敗。最後錯誤: {last_error}")
            
            # 立即驗證文件是否真的存在
            import time
            time.sleep(0.2)  # 等待文件系統更新
            
            # 檢查多個可能的路徑
            check_paths = [
                Path(save_path),
                Path(save_path).resolve(),
            ]
            
            file_exists = False
            actual_path = None
            for check_path in check_paths:
                if check_path.exists():
                    file_exists = True
                    actual_path = check_path
                    break
            
            if not file_exists:
                error_msg = f"[ERROR] 文件保存後不存在！保存路徑: {save_path}"
                print(error_msg)
                print(f"[DEBUG] 父目錄是否存在: {Path(save_path).parent.exists()}")
                print(f"[DEBUG] 父目錄: {Path(save_path).parent}")
                print(f"[DEBUG] 父目錄是否可寫: {os.access(Path(save_path).parent, os.W_OK)}")
                # 列出父目錄中的所有文件
                try:
                    files_in_parent = list(Path(save_path).parent.iterdir())
                    print(f"[DEBUG] 父目錄中的文件: {files_in_parent}")
                except:
                    pass
                raise Exception(error_msg)
            
            # 使用實際存在的路徑
            save_path = str(actual_path)
            print(f"[DEBUG] 文件確認存在於: {save_path}")
            
            file_size = Path(save_path).stat().st_size
            print(f"[DEBUG] Step 3 成功: 文件已保存")
            print(f"[DEBUG] 文件大小: {file_size} bytes")
            
            # 驗證文件大小是否合理（至少應該有幾 KB）
            if file_size < 1000:  # 小於 1KB 可能不正常
                print(f"[WARNING] 文件大小異常小: {file_size} bytes，可能保存不完整")
                try:
                    st.warning(f"⚠️ 文件大小異常小 ({file_size} bytes)，可能保存不完整")
                except:
                    pass
            
        except Exception as save_error:
            error_msg = f"[ERROR] Failed to save PPTX file: {str(save_error)}"
            print(error_msg)
            import traceback
            full_traceback = traceback.format_exc()
            print(full_traceback)
            try:
                st.error(f"❌ 文件保存失敗: {str(save_error)}")
                with st.expander("詳細錯誤信息", expanded=True):
                    st.code(full_traceback)
            except:
                pass
            raise Exception(error_msg) from save_error
        
        # 更新會話活動時間（用於會話清理）
        try:
            print("[DEBUG] Step 4: 更新會話活動時間...")
            update_session_activity()
            print("[DEBUG] Step 4 成功")
        except Exception as activity_error:
            # 更新活動時間失敗不影響主流程，只記錄警告
            print(f"[WARNING] Failed to update session activity: {str(activity_error)}")
            try:
                st.warning(f"⚠️ 更新會話活動時間失敗: {str(activity_error)}")
            except:
                pass
        
        # 保存路徑到 session_state（供其他模組讀取）
        try:
            print("[DEBUG] Step 5: 保存路徑到 session_state...")
            st.session_state["tcfd_report_file"] = str(output_path)
            print(f"[DEBUG] Step 5 成功: {st.session_state.get('tcfd_report_file')}")
        except Exception as state_error:
            # session_state 保存失敗不影響主流程，只記錄警告
            print(f"[WARNING] Failed to save path to session_state: {str(state_error)}")
            try:
                st.warning(f"⚠️ 保存路徑到 session_state 失敗: {str(state_error)}")
            except:
                pass
        
        # 最終驗證：確保文件真的存在且有效
        save_path = str(output_path)
        final_path = Path(save_path)
        
        if not final_path.exists():
            error_msg = f"[ERROR] 最終驗證失敗：文件不存在！路徑: {save_path}"
            print(error_msg)
            try:
                st.error(f"❌ 文件保存驗證失敗：文件不存在於 `{save_path}`")
                st.info("💡 文件可能沒有真正保存，請檢查權限和路徑")
            except:
                pass
            raise Exception(error_msg)
        
        file_size = final_path.stat().st_size
        print("[DEBUG] ========== TCFD 文件保存完成 ==========")
        print(f"[DEBUG] 最終輸出路徑: {save_path}")
        print(f"[DEBUG] 最終輸出路徑 (absolute): {final_path.resolve()}")
        print(f"[DEBUG] 文件大小: {file_size} bytes")
        
        # 在 UI 中顯示保存成功訊息（只有文件真的存在時才顯示）
        try:
            file_size_kb = file_size / 1024
            st.success(f"✅ **文件已成功保存並驗證！**\n\n"
                      f"📁 **路徑**: `{save_path}`\n\n"
                      f"📊 **文件大小**: {file_size_kb:.2f} KB ({file_size:,} bytes)\n\n"
                      f"📄 **Slides 數量**: {len(prs.slides)} 頁")
            print(f"[SUCCESS] 文件保存成功並在 UI 中顯示: {save_path}")
        except Exception as display_error:
            print(f"[WARNING] 無法在 UI 中顯示成功訊息: {display_error}")
            # 即使顯示失敗，也不影響返回路徑
        
        return final_path
        
    except Exception as e:
        error_msg = f"[ERROR] Error generating combined PPTX: {str(e)}"
        print(error_msg)
        import traceback
        full_traceback = traceback.format_exc()
        print("=" * 60)
        print("Full traceback:")
        print(full_traceback)
        print("=" * 60)
        # 將錯誤信息也記錄到 stderr，確保 Streamlit 能看到
        import sys
        sys.stderr.write(error_msg + "\n")
        sys.stderr.write(full_traceback + "\n")
        
        # 強制在 Streamlit 中顯示錯誤（如果可用）
        # 注意：這裡的 st 可能不在正確的上下文中，所以讓上層處理
        # 但我們確保錯誤信息被完整記錄
        print("[ERROR] ========== 錯誤已記錄，請查看 UI 中的錯誤顯示 ==========")
        
        # 重新拋出異常，讓上層捕獲並顯示
        raise Exception(f"生成 TCFD 報告失敗: {str(e)}\n\n詳細信息請查看終端輸出") from e

