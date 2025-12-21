from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

# ================= 🛠️ 基礎工具函數 (保持不變) =================
def set_cell_bg(cell, hex_color):
    if not hex_color: return
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def set_text(cell, text, font_size=10, is_bold=False, color='000000', align='center'):
    if not cell.text_frame.paragraphs:
        cell.text_frame.add_paragraph()
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT
    p.text = "" 
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = is_bold
    run.font.name = 'Arial'
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def remove_all_borders(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = OxmlElement(f'a:{edge}')
        ln.set('w', '0')
        noFill = OxmlElement('a:noFill')
        ln.append(noFill)
        tcPr.append(ln)

def init_zebra_table(slide, rows=4, cols=6):
    left, top = Inches(0.5), Inches(1.0)
    width, height = Inches(12.0), Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 欄寬設定: 強調 Action 與 Financial Logic
    col_widths = [1.2, 1.5, 2.0, 2.8, 2.5, 2.0]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
        
    for r in range(rows):
        for c in range(cols):
            remove_all_borders(table.cell(r, c))
    return table

# 色彩定義 (專業冷色調)
COLOR_BG_WHITE = 'FFFFFF'
COLOR_BG_HEADER = 'EFEFEF'
COLOR_BG_STRIPE = 'F7F7F7'
COLOR_TEXT_SUB = '333333'

# ================= 📝 Table 6: Systemic Risk Control (原 Governance) =================
def create_slide_6_risk_control(prs, data_lines=None):
    # 動態查找空白 layout
    blank_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        layout_name_lower = layout.name.lower()
        if 'blank' in layout_name_lower or 'empty' in layout_name_lower:
            blank_layout = layout
            break
    if blank_layout is None and len(prs.slide_layouts) > 6:
        blank_layout = prs.slide_layouts[6]
    elif blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    
    slide = prs.slides.add_slide(blank_layout)
    table = init_zebra_table(slide)

    # 1. 主標題 (完全移除 Governance 字眼)
    cell = table.cell(0, 0); cell.merge(table.cell(0, 2))
    set_text(cell, "Systemic Risk Control", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)
    cell = table.cell(0, 3); cell.merge(table.cell(0, 5))
    set_text(cell, "Infrastructure & Assurance", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)

    # 2. 欄位標題 (強調控制與防護)
    headers = ['Control Area', 'External Driver /\nRequirement', 'System Gap /\nExposure', 'Mitigation Protocol\n(Soft Infrastructure)', 'Liability Avoidance /\nCost Benefit', 'Budget\nAllocation']
    for i, h in enumerate(headers):
        set_text(table.cell(1, i), h, 10, True, COLOR_TEXT_SUB)
        set_cell_bg(table.cell(1, i), COLOR_BG_HEADER)

    # --- Row 2: Data Integrity (Scope 3) ---
    set_text(table.cell(2, 0), "Data Integrity\n(Scope 3)", 9, True, COLOR_TEXT_SUB)
    set_text(table.cell(2, 1), "Mandatory 3rd-Party\nVerification", 9, False, COLOR_TEXT_SUB)
    set_text(table.cell(2, 2), "Unverified Upstream\nData", 9, False, COLOR_TEXT_SUB) 
    # 剩下的由 LLM 填寫
    for c in range(0, 6): set_cell_bg(table.cell(2, c), COLOR_BG_WHITE)

    # --- Row 3: Supply Chain Integrity (Due Diligence) ---
    set_text(table.cell(3, 0), "Supply Chain\nIntegrity", 9, True, COLOR_TEXT_SUB)
    set_text(table.cell(3, 1), "Traceability &\nDue Diligence", 9, False, COLOR_TEXT_SUB)
    set_text(table.cell(3, 2), "Tier-2 Visibility\nGap", 9, False, COLOR_TEXT_SUB)
    for c in range(3, 6): 
        set_text(table.cell(3, c), "", 9)
        table.cell(3, c).vertical_anchor = MSO_ANCHOR.TOP
    for c in range(0, 6): set_cell_bg(table.cell(3, c), COLOR_BG_STRIPE)
    
    # 填充 LLM 返回的數據（列 3-5: Mitigation Protocol, Liability Avoidance, Budget）
    if data_lines:
        data_rows = [2, 3]
        for idx, data_line in enumerate(data_lines[:2]):
            if idx >= len(data_rows):
                break
            row_idx = data_rows[idx]
            parts = [p.strip() for p in data_line.split('|||')]
            # table06 的列：3=Mitigation Protocol, 4=Liability Avoidance, 5=Budget
            if len(parts) >= 1 and parts[0]:
                desc_parts = parts[0].split(';', 1)
                mitigation = desc_parts[1].strip() if len(desc_parts) > 1 else desc_parts[0].strip()
                if mitigation:
                    set_text(table.cell(row_idx, 3), mitigation, 9)
            if len(parts) >= 2 and parts[1]:
                set_text(table.cell(row_idx, 4), parts[1].strip(), 9)  # Liability Avoidance
            if len(parts) >= 3 and parts[2]:
                set_text(table.cell(row_idx, 5), parts[2].strip(), 9)  # Budget


# ================= 📝 Table 7: Operational Resilience (原 Social) =================
def create_slide_7_resilience(prs, data_lines=None):
    # 動態查找空白 layout
    blank_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        layout_name_lower = layout.name.lower()
        if 'blank' in layout_name_lower or 'empty' in layout_name_lower:
            blank_layout = layout
            break
    if blank_layout is None and len(prs.slide_layouts) > 6:
        blank_layout = prs.slide_layouts[6]
    elif blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    
    slide = prs.slides.add_slide(blank_layout)
    table = init_zebra_table(slide)

    # 1. 主標題 (完全移除 Social 字眼，使用 IPCC 術語)
    cell = table.cell(0, 0); cell.merge(table.cell(0, 2))
    set_text(cell, "Operational Resilience", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)
    cell = table.cell(0, 3); cell.merge(table.cell(0, 5))
    set_text(cell, "Adaptive Capacity (Human & Supply)", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)

    # 2. 欄位標題 (強調運作連續性與產能)
    headers = ['Resilience Unit', 'Physical/Tech\nStressor', 'Operational Impact\n(Downtime Risk)', 'Adaptation Strategy\n(Capacity Building)', 'Continuity Benefit\n(ROI)', 'Budget\nAllocation']
    for i, h in enumerate(headers):
        set_text(table.cell(1, i), h, 10, True, COLOR_TEXT_SUB)
        set_cell_bg(table.cell(1, i), COLOR_BG_HEADER)

    # --- Row 2: Workforce (Heat/Safety) ---
    set_text(table.cell(2, 0), "Workforce\nAdaptation", 9, True, COLOR_TEXT_SUB)
    set_text(table.cell(2, 1), "Thermal Stress /\nNew Process Risks", 9, False, COLOR_TEXT_SUB)
    set_text(table.cell(2, 2), "Productivity Loss\n(-15% Forecast)", 9, False, COLOR_TEXT_SUB) 
    for c in range(0, 6): set_cell_bg(table.cell(2, c), COLOR_BG_WHITE)

    # --- Row 3: Value Chain (Water/Resource) ---
    set_text(table.cell(3, 0), "Value Chain\nSecurity", 9, True, COLOR_TEXT_SUB)
    set_text(table.cell(3, 1), "Resource Competition\n(Water/Power)", 9, False, COLOR_TEXT_SUB)
    set_text(table.cell(3, 2), "License to Operate\nRevocation", 9, False, COLOR_TEXT_SUB)
    for c in range(3, 6): 
        set_text(table.cell(3, c), "", 9)
        table.cell(3, c).vertical_anchor = MSO_ANCHOR.TOP
    for c in range(0, 6): set_cell_bg(table.cell(3, c), COLOR_BG_STRIPE)
    
    # 填充 LLM 返回的數據（列 3-5: Adaptation Strategy, Continuity Benefit, Budget）
    if data_lines:
        data_rows = [2, 3]
        for idx, data_line in enumerate(data_lines[:2]):
            if idx >= len(data_rows):
                break
            row_idx = data_rows[idx]
            parts = [p.strip() for p in data_line.split('|||')]
            # table07 的列：3=Adaptation Strategy, 4=Continuity Benefit, 5=Budget
            if len(parts) >= 1 and parts[0]:
                desc_parts = parts[0].split(';', 1)
                adaptation = desc_parts[1].strip() if len(desc_parts) > 1 else desc_parts[0].strip()
                if adaptation:
                    set_text(table.cell(row_idx, 3), adaptation, 9)
            if len(parts) >= 2 and parts[1]:
                set_text(table.cell(row_idx, 4), parts[1].strip(), 9)  # Continuity Benefit
            if len(parts) >= 3 and parts[2]:
                set_text(table.cell(row_idx, 5), parts[2].strip(), 9)  # Budget


def generate_table_06(data_lines=None, filename=None, prs=None):
    # 如果提供了 prs，直接添加到主 prs；否則創建獨立文件（向後兼容）
    if prs is not None:
        create_slide_6_risk_control(prs, data_lines=data_lines)
        return None  # 已添加到主 prs，不需要返回文件名
    else:
        if filename is None:
            filename = 'TCFD_table06_systemic_risk.pptx'
        prs = Presentation()
        create_slide_6_risk_control(prs, data_lines=data_lines)
        prs.save(filename)
        return filename

def generate_table_07(data_lines=None, filename=None, prs=None):
    # 如果提供了 prs，直接添加到主 prs；否則創建獨立文件（向後兼容）
    if prs is not None:
        create_slide_7_resilience(prs, data_lines=data_lines)
        return None  # 已添加到主 prs，不需要返回文件名
    else:
        if filename is None:
            filename = 'TCFD_table07_operational_resilience.pptx'
        prs = Presentation()
        create_slide_7_resilience(prs, data_lines=data_lines)
        prs.save(filename)
        return filename

if __name__ == "__main__":
    prs = Presentation()
    # 建立 Table 6 (Systemic Risk Control)
    create_slide_6_risk_control(prs)
    # 建立 Table 7 (Operational Resilience)
    create_slide_7_resilience(prs)
    
    filename = "tcfd_tables_6_7_strategic.pptx"
    prs.save(filename)
    print(f"✅ Table 6 & 7 生成完成 (已更名): {filename}")
    
    try:
        from google.colab import files
        files.download(filename)
    except:
        pass