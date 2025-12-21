from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

# ================= 🛠️ 底層繪圖工具 =================
def set_cell_bg(cell, hex_color):
    if not hex_color: return
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def set_text(cell, text, font_size=10, is_bold=False, color='000000', align='center'):
    if not cell.text_frame.paragraphs:
        cell.text_frame.add_paragraph()
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT
    p.text = "" 
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = is_bold
    run.font.name = 'Arial'
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def remove_all_borders(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = OxmlElement(f'a:{edge}')
        ln.set('w', '0')
        noFill = OxmlElement('a:noFill')
        ln.append(noFill)
        tcPr.append(ln)

def init_zebra_table(slide, rows=6, cols=6):
    left, top = Inches(0.5), Inches(0.8)
    width, height = Inches(12.0), Inches(5.8)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 欄寬設定 (確保標題斷行正確)
    col_widths = [0.9, 1.3, 0.9, 2.7, 2.7, 2.1]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
        
    # 移除預設格線
    for r in range(rows):
        for c in range(cols):
            remove_all_borders(table.cell(r, c))
            
    return table

# 色彩定義
COLOR_BG_WHITE = 'FFFFFF'
COLOR_BG_HEADER = 'EFEFEF'
COLOR_BG_STRIPE = 'F7F7F7'
COLOR_TEXT_SUB = '333333'

# ================= 📝 修正後的 Transformation Risk 表格 =================
def create_slide_transformation_corrected(prs=None, output_filename=None, data_lines=None):
    # 如果提供了 prs，直接使用；否則創建新的（向後兼容）
    if prs is None:
        prs = Presentation()
        output_mode = True
    else:
        output_mode = False
    
    # 動態查找空白 layout
    blank_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        layout_name_lower = layout.name.lower()
        if 'blank' in layout_name_lower or 'empty' in layout_name_lower:
            blank_layout = layout
            break
    if blank_layout is None and len(prs.slide_layouts) > 6:
        blank_layout = prs.slide_layouts[6]
    elif blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    
    slide = prs.slides.add_slide(blank_layout)
    table = init_zebra_table(slide)

    # 1. 主標題 (Climate-Related Risks / Financial Impacts)
    cell = table.cell(0, 0); cell.merge(table.cell(0, 2))
    set_text(cell, "Climate-Related Risks", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)
    
    cell = table.cell(0, 3); cell.merge(table.cell(0, 5))
    set_text(cell, "Financial Impacts", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)

    # 2. 欄位標題 (包含您指定的強制斷行)
    headers = ['Type', 'Climate\nChange\nRelated Factor', 'Impact\nPeriod', 'Description of Content', 'Potential Financial Impact', 'Adaption & Response']
    for i, h in enumerate(headers):
        set_text(table.cell(1, i), h, 10, True, COLOR_TEXT_SUB)
        set_cell_bg(table.cell(1, i), COLOR_BG_HEADER)

    # 3. 內容填充
    
    # --- Row 2 (Policy) ---
    # Type: Transformation Risk
    set_text(table.cell(2, 0), "Transformation\nRisk", 9, True, COLOR_TEXT_SUB)
    
    # Factor: Policy
    set_text(table.cell(2, 1), "Policy and\nRegulation", 9, False, COLOR_TEXT_SUB)
    
    # Period: Short-medium (第一列)
    set_text(table.cell(2, 2), "Short-term\nand\nmedium-term", 9, False, COLOR_TEXT_SUB)
    
    # 其他欄位留白 (Row 2 是白底)
    for c in range(3, 6): set_cell_bg(table.cell(2, c), COLOR_BG_WHITE)
    # 補上 Col 0~2 的白底
    for c in range(0, 3): set_cell_bg(table.cell(2, c), COLOR_BG_WHITE)


    # --- Row 3 (Green product) ---
    # Type: 與上方合併 (Transformation Risk)
    table.cell(2, 0).merge(table.cell(3, 0)) # 合併後會維持 Row 2 的白色背景

    # Factor: Green product
    set_text(table.cell(3, 1), "Green product\nand technology", 9, False, COLOR_TEXT_SUB)
    set_cell_bg(table.cell(3, 1), COLOR_BG_STRIPE) # 灰底

    # Period: Short-medium (第二列 - 這裡原本是 Long-term，現已修正！)
    set_text(table.cell(3, 2), "Short-term\nand\nmedium-term", 9, False, COLOR_TEXT_SUB)
    set_cell_bg(table.cell(3, 2), COLOR_BG_STRIPE) # 灰底

    # 其他欄位留白 (Row 3 是灰底)
    for c in range(3, 6): 
        set_text(table.cell(3, c), "", 9)
        table.cell(3, c).vertical_anchor = MSO_ANCHOR.TOP
        set_cell_bg(table.cell(3, c), COLOR_BG_STRIPE)
    
    # --- Row 4 (預設格式，白底) ---
    table.cell(2, 0).merge(table.cell(4, 0))  # 繼續合併 Type 列
    set_text(table.cell(4, 1), "Market\nDisruption", 9, False, COLOR_TEXT_SUB)
    set_text(table.cell(4, 2), "Medium-term", 9, False, COLOR_TEXT_SUB)
    for c in range(3, 6):
        set_cell_bg(table.cell(4, c), COLOR_BG_WHITE)
    for c in range(0, 3):
        set_cell_bg(table.cell(4, c), COLOR_BG_WHITE)
    
    # --- Row 5 (預設格式，灰底) ---
    table.cell(2, 0).merge(table.cell(5, 0))  # 繼續合併 Type 列
    set_text(table.cell(5, 1), "Reputation\nRisk", 9, False, COLOR_TEXT_SUB)
    set_text(table.cell(5, 2), "Long-term", 9, False, COLOR_TEXT_SUB)
    for c in range(3, 6):
        set_text(table.cell(5, c), "", 9)
        table.cell(5, c).vertical_anchor = MSO_ANCHOR.TOP
        set_cell_bg(table.cell(5, c), COLOR_BG_STRIPE)
    for c in range(0, 3):
        set_cell_bg(table.cell(5, c), COLOR_BG_STRIPE)
    
    # 填充 LLM 返回的數據（如果有）
    if data_lines:
        # 解析 data_lines（格式：Description ||| Financial Impact ||| Adaptation）
        # 表格結構：Row 2-5 的第 3-5 列需要填充（支持最多4行數據）
        data_rows = [2, 3, 4, 5]  # 需要填充數據的行
        
        for idx, data_line in enumerate(data_lines[:4]):  # 最多4行數據
            if idx >= len(data_rows):
                break
            
            row_idx = data_rows[idx]
            
            # 解析 ||| 分隔的數據
            parts = [p.strip() for p in data_line.split('|||')]
            
            # 處理第一部分（可能包含分號分隔的標題和描述）
            description = ""
            if len(parts) >= 1 and parts[0]:
                desc_parts = parts[0].split(';', 1)  # 分號分隔標題和描述
                if len(desc_parts) > 1:
                    description = desc_parts[1].strip()  # 取描述部分
                else:
                    description = desc_parts[0].strip()
            
            # 填充到對應的列（Description, Financial Impact, Adaptation）
            # 使用較小的字體以容納更多文字
            if description:
                set_text(table.cell(row_idx, 3), description, 8)  # 從 9 降到 8
            if len(parts) >= 2 and parts[1]:
                set_text(table.cell(row_idx, 4), parts[1].strip(), 8)  # Financial Impact
            if len(parts) >= 3 and parts[2]:
                set_text(table.cell(row_idx, 5), parts[2].strip(), 8)  # Adaptation

    # 只有在獨立模式下才保存
    if output_mode and output_filename:
        prs.save(output_filename)
        print(f"✅ Transformation Risk 表格修正完成: {output_filename}")


if __name__ == "__main__":
    create_slide_transformation_corrected("tcfd_slide_transformation_corrected.pptx")
    try:
        from google.colab import files
        files.download('tcfd_slide_transformation_corrected.pptx')
    except:
        pass

def generate_table_01(data_lines=None, filename=None, prs=None):
    # 如果提供了 prs，直接添加到主 prs；否則創建獨立文件（向後兼容）
    if prs is not None:
        create_slide_transformation_corrected(prs=prs, data_lines=data_lines)
        return None  # 已添加到主 prs，不需要返回文件名
    else:
        if filename is None:
            filename = 'TCFD_table01_transformation_risks.pptx'
        create_slide_transformation_corrected(output_filename=filename, data_lines=data_lines)
        return filename
