from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

# ================= 🛠️ 基礎工具函數 =================
def set_cell_bg(cell, hex_color):
    if not hex_color: return
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def set_text(cell, text, font_size=10, is_bold=False, color='000000', align='center'):
    if not cell.text_frame.paragraphs:
        cell.text_frame.add_paragraph()
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT
    p.text = "" 
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = is_bold
    run.font.name = 'Arial'
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def remove_all_borders(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = OxmlElement(f'a:{edge}')
        ln.set('w', '0')
        noFill = OxmlElement('a:noFill')
        ln.append(noFill)
        tcPr.append(ln)

def init_zebra_table(slide, rows=4, cols=6):
    left, top = Inches(0.5), Inches(1.0)
    width, height = Inches(12.0), Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 欄寬設定 (指標表通常需要較寬的"Progress"欄位，稍微調整)
    # Col 0: Type, Col 1: Metric, Col 2: Year, Col 3: Progress...
    col_widths = [1.2, 1.2, 1.0, 2.5, 2.5, 2.0]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
        
    # 移除所有預設格線
    for r in range(rows):
        for c in range(cols):
            remove_all_borders(table.cell(r, c))
            
    return table

# 色彩定義
COLOR_BG_WHITE = 'FFFFFF'
COLOR_BG_HEADER = 'EFEFEF'
COLOR_BG_STRIPE = 'F7F7F7'
COLOR_TEXT_SUB = '333333'

# ================= 📝 Table 5: Metrics & Targets 生成邏輯 =================
def create_slide_5_metrics(prs=None, output_filename=None, data_lines=None):
    # 如果提供了 prs，直接使用；否則創建新的（向後兼容）
    if prs is None:
        prs = Presentation()
        output_mode = True
    else:
        output_mode = False
    
    # 動態查找空白 layout
    blank_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        layout_name_lower = layout.name.lower()
        if 'blank' in layout_name_lower or 'empty' in layout_name_lower:
            blank_layout = layout
            break
    if blank_layout is None and len(prs.slide_layouts) > 6:
        blank_layout = prs.slide_layouts[6]
    elif blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    
    slide = prs.slides.add_slide(blank_layout)
    table = init_zebra_table(slide)

    # 1. 主標題
    cell = table.cell(0, 0); cell.merge(table.cell(0, 2))
    set_text(cell, "Metrics & Targets", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)
    
    cell = table.cell(0, 3); cell.merge(table.cell(0, 5))
    set_text(cell, "Performance Indicators", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)

    # 2. 欄位標題 (切換為指標專用欄位)
    # 注意：這裡使用了更適合量化數據的標題
    headers = ['Type', 'Metric\nCategory', 'Target\nTimeframe', 'Current Progress\n& Status', 'Financial\nLinkage', 'Action Plan']
    for i, h in enumerate(headers):
        set_text(table.cell(1, i), h, 10, True, COLOR_TEXT_SUB)
        set_cell_bg(table.cell(1, i), COLOR_BG_HEADER)

    # 3. 內容填充
    
    # --- Row 2: GHG Emissions (白底) ---
    # Type
    set_text(table.cell(2, 0), "GHG Emissions\n(Scope 1, 2, 3)", 9, True, COLOR_TEXT_SUB)
    
    # Metric (範例)
    set_text(table.cell(2, 1), "Emission Reduction\nTargets", 9, False, COLOR_TEXT_SUB)
    
    # Timeframe: 短中期 (依您指示)
    set_text(table.cell(2, 2), "Short-term\nand\nmedium-term", 9, False, COLOR_TEXT_SUB)
    
    # 設為白底
    for c in range(0, 6): set_cell_bg(table.cell(2, c), COLOR_BG_WHITE)


    # --- Row 3: Climate-Related Targets (灰底) ---
    # Type
    set_text(table.cell(3, 0), "Climate-Related\nTargets", 9, True, COLOR_TEXT_SUB)
    
    # Metric (範例: 水與廢棄物)
    set_text(table.cell(3, 1), "Water, Waste &\nGreen Revenue", 9, False, COLOR_TEXT_SUB)
    
    # Timeframe: 短中期 (依您指示)
    set_text(table.cell(3, 2), "Short-term\nand\nmedium-term", 9, False, COLOR_TEXT_SUB)
    
    # 其他空白欄位靠上對齊
    for c in range(3, 6): 
        set_text(table.cell(3, c), "", 9)
        table.cell(3, c).vertical_anchor = MSO_ANCHOR.TOP

    # 設為灰底
    for c in range(0, 6): set_cell_bg(table.cell(3, c), COLOR_BG_STRIPE)
    
    # 填充 LLM 返回的數據
    if data_lines:
        data_rows = [2, 3]
        for idx, data_line in enumerate(data_lines[:2]):
            if idx >= len(data_rows):
                break
            row_idx = data_rows[idx]
            parts = [p.strip() for p in data_line.split('|||')]
            description = ""
            if len(parts) >= 1 and parts[0]:
                desc_parts = parts[0].split(';', 1)
                description = desc_parts[1].strip() if len(desc_parts) > 1 else desc_parts[0].strip()
            if description:
                set_text(table.cell(row_idx, 3), description, 9)
            if len(parts) >= 2 and parts[1]:
                set_text(table.cell(row_idx, 4), parts[1].strip(), 9)
            if len(parts) >= 3 and parts[2]:
                set_text(table.cell(row_idx, 5), parts[2].strip(), 9)

    # 只有在獨立模式下才保存
    # 只有在獨立模式下才保存
    if output_mode and output_filename:
        prs.save(output_filename)
        print(f"✅ Table 5 (Metrics & Targets) 生成完成: {output_filename}")

def generate_table_05(data_lines=None, filename=None, prs=None):
    # 如果提供了 prs，直接添加到主 prs；否則創建獨立文件（向後兼容）
    if prs is not None:
        create_slide_5_metrics(prs=prs, data_lines=data_lines)
        return None  # 已添加到主 prs，不需要返回文件名
    else:
        if filename is None:
            filename = 'TCFD_table05_metrics_targets.pptx'
        create_slide_5_metrics(output_filename=filename, data_lines=data_lines)
        return filename

if __name__ == "__main__":
    create_slide_5_metrics("tcfd_table_5_metrics.pptx")
    try:
        from google.colab import files
        files.download('tcfd_table_5_metrics.pptx')
    except:
        pass