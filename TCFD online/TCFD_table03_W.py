from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

# ================= 🛠️ 基礎工具函數 =================
def set_cell_bg(cell, hex_color):
    if not hex_color: return
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)

def set_text(cell, text, font_size=10, is_bold=False, color='000000', align='center'):
    if not cell.text_frame.paragraphs:
        cell.text_frame.add_paragraph()
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT
    p.text = "" 
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = is_bold
    run.font.name = 'Arial'
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def remove_all_borders(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = OxmlElement(f'a:{edge}')
        ln.set('w', '0')
        noFill = OxmlElement('a:noFill')
        ln.append(noFill)
        tcPr.append(ln)

def init_zebra_table(slide, rows=4, cols=6):
    left, top = Inches(0.5), Inches(1.0)
    width, height = Inches(12.0), Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 欄寬設定 (確保標題斷行正確)
    col_widths = [0.9, 1.3, 0.9, 2.7, 2.7, 2.1]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
        
    # 移除所有預設格線
    for r in range(rows):
        for c in range(cols):
            remove_all_borders(table.cell(r, c))
            
    return table

# 色彩定義
COLOR_BG_WHITE = 'FFFFFF'
COLOR_BG_HEADER = 'EFEFEF'
COLOR_BG_STRIPE = 'F7F7F7'
COLOR_TEXT_SUB = '333333'

# ================= 📝 Table 3: Resource & Energy 生成邏輯 =================
def create_slide_3_resource_energy(output_filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = init_zebra_table(slide)

    # 1. 主標題
    cell = table.cell(0, 0); cell.merge(table.cell(0, 2))
    set_text(cell, "Climate Opportunities", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)
    
    cell = table.cell(0, 3); cell.merge(table.cell(0, 5))
    set_text(cell, "Financial Benefits", 16, True); set_cell_bg(cell, COLOR_BG_WHITE)

    # 2. 欄位標題
    headers = ['Type', 'Climate\nChange\nRelated Factor', 'Impact\nPeriod', 'Description of Content', 'Potential Financial Impact', 'Adaption & Response']
    for i, h in enumerate(headers):
        set_text(table.cell(1, i), h, 10, True, COLOR_TEXT_SUB)
        set_cell_bg(table.cell(1, i), COLOR_BG_HEADER)

    # 3. 內容填充
    
    # --- Row 2: Resource Efficiency (白底) ---
    # Type
    set_text(table.cell(2, 0), "Resource\nEfficiency", 9, True, COLOR_TEXT_SUB)
    
    # Factor (範例: 製程優化)
    set_text(table.cell(2, 1), "Production process\nefficiency", 9, False, COLOR_TEXT_SUB)
    
    # Period: 短中期 (依您指示)
    set_text(table.cell(2, 2), "Short-term\nand\nmedium-term", 9, False, COLOR_TEXT_SUB)
    
    # 設為白底
    for c in range(0, 6): set_cell_bg(table.cell(2, c), COLOR_BG_WHITE)


    # --- Row 3: Energy Source (灰底) ---
    # Type
    set_text(table.cell(3, 0), "Energy\nSource", 9, True, COLOR_TEXT_SUB)
    
    # Factor (範例: 低碳能源)
    set_text(table.cell(3, 1), "Use of lower-emission\nsources of energy", 9, False, COLOR_TEXT_SUB)
    
    # Period: 短中期 (依您指示)
    set_text(table.cell(3, 2), "Short-term\nand\nmedium-term", 9, False, COLOR_TEXT_SUB)
    
    # 其他空白欄位靠上對齊
    for c in range(3, 6): 
        set_text(table.cell(3, c), "", 9)
        table.cell(3, c).vertical_anchor = MSO_ANCHOR.TOP

    # 設為灰底
    for c in range(0, 6): set_cell_bg(table.cell(3, c), COLOR_BG_STRIPE)

    # 存檔
    prs.save(output_filename)
    print(f"✅ Table 3 (Resource & Energy) 生成完成: {output_filename}")

if __name__ == "__main__":
    create_slide_3_resource_energy("tcfd_table_3_resource_energy.pptx")
    try:
        from google.colab import files
        files.download('tcfd_table_3_resource_energy.pptx')
    except:
        pass