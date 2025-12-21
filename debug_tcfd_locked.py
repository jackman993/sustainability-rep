"""
測試文件被鎖定時的行為
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

print("=" * 60)
print("測試文件鎖定情況")
print("=" * 60)

from pptx import Presentation

output_dir = Path(__file__).parent / "shared" / "engine" / "tcfd" / "output"
output_path = output_dir / "TCFD_table.pptx"

print(f"\n1. 檢查輸出文件：")
print(f"   文件: {output_path}")
print(f"   存在: {output_path.exists()}")

if output_path.exists():
    print(f"\n2. 嘗試打開文件（模擬 PowerPoint 鎖定）...")
    try:
        # 嘗試以讀寫模式打開文件
        with open(output_path, 'r+b') as f:
            print(f"   ✅ 文件可以打開（未鎖定）")
            # 讀取一點內容
            f.read(1)
    except PermissionError as e:
        print(f"   ❌ 文件被鎖定: {e}")
    except Exception as e:
        print(f"   ⚠️ 其他錯誤: {e}")

print(f"\n3. 嘗試刪除文件...")
try:
    if output_path.exists():
        output_path.unlink()
        print(f"   ✅ 文件刪除成功")
    else:
        print(f"   ℹ️ 文件不存在，無需刪除")
except PermissionError as e:
    print(f"   ❌ 無法刪除（文件被鎖定）: {e}")
    print(f"   這會觸發代碼中的時間戳重命名邏輯")
except Exception as e:
    print(f"   ⚠️ 其他錯誤: {e}")

print(f"\n4. 嘗試保存新文件...")
try:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    
    # 如果文件存在且無法刪除，使用時間戳
    if output_path.exists():
        try:
            output_path.unlink()
        except PermissionError:
            from datetime import datetime
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = output_dir / f"TCFD_table_{timestamp}.pptx"
            print(f"   使用新文件名: {output_path.name}")
    
    prs.save(str(output_path))
    print(f"   ✅ 文件保存成功: {output_path}")
    print(f"   文件大小: {output_path.stat().st_size} bytes")
    
except PermissionError as e:
    print(f"   ❌ 保存失敗（權限錯誤）: {e}")
    print(f"   這可能是因為文件被 PowerPoint 鎖定")
except Exception as e:
    print(f"   ❌ 保存失敗: {type(e).__name__}: {e}")
    import traceback
    traceback.print_exc()

